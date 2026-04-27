"""엔지니어링 심층 자료 (KO). 출처: 이슈 #1 (SSOT 아키텍처), #4 (한국어 NLP / GEO 기술스택)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from build_pitch_deck import (
    NAVY, INK, TEAL, MUTED, GOLD, DIM,
    W, H, add_bg, bar,
)

KFONT = "Pretendard"


def tx(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=KFONT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return box


def chip(slide, x, y, label, color=TEAL):
    h = Inches(0.32)
    w = Inches(0.07 * len(label) + 0.4)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.5; s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = DIM
    tf = s.text_frame
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = KFONT; r.font.size = Pt(10)
    r.font.bold = True; r.font.color.rgb = color


def header(slide, eyebrow, title):
    tx(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
       eyebrow, size=11, bold=True, color=TEAL)
    bar(slide, Inches(0.6), Inches(0.95), w=Inches(0.5), h=Inches(0.06))
    tx(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.9),
       title, size=30, bold=True, color=INK)


def footer(slide, page, total, label="aeogeo · yogaman.club"):
    tx(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
       label, size=10, color=MUTED)
    tx(slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3),
       f"{page} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, fill=DIM):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.05; s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s


def kv_row(slide, x, y, w, key, val, key_w=Inches(2.0)):
    tx(slide, x, y, key_w, Inches(0.32), key, size=11, bold=True, color=TEAL)
    tx(slide, x + key_w, y, w - key_w, Inches(0.32), val, size=11, color=INK)


def bullet(slide, x, y, w, h, title, body, accent=TEAL):
    bar(slide, x, y + Inches(0.04), w=Inches(0.06), h=Inches(0.42), color=accent)
    tx(slide, x + Inches(0.18), y, w - Inches(0.2), Inches(0.4),
       title, size=14, bold=True, color=INK)
    tx(slide, x + Inches(0.18), y + Inches(0.45), w - Inches(0.2), h - Inches(0.45),
       body, size=11, color=MUTED)


def code_block(slide, x, y, w, h, text):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.line.fill.background(); s.fill.solid(); s.fill.fore_color.rgb = DIM
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = tf.margin_bottom = Inches(0.14)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = line
        r.font.name = "D2Coding"; r.font.size = Pt(10)
        r.font.color.rgb = TEAL if line.startswith(("→", "│", "├", "└", "┌", "┐")) else INK


# ---------- slides ----------
def s_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               Inches(0.18), H)
    block.line.fill.background(); block.fill.solid()
    block.fill.fore_color.rgb = TEAL
    tx(s, Inches(0.8), Inches(0.6), Inches(6), Inches(0.5),
       "AEOGEO · 엔지니어링", size=14, bold=True, color=TEAL)
    tx(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.4),
       "엔지니어링 심층 자료", size=54, bold=True, color=INK)
    tx(s, Inches(0.8), Inches(3.7), Inches(11.5), Inches(1.6),
       "165쪽 도메인 책을 한국어 RAG API로 변환해\n"
       "AI 어시스턴트가 실제로 인용하게 만드는 방법.",
       size=20, color=MUTED)
    chip(s, Inches(0.8), Inches(5.4), "FastAPI · Ollama · Tesseract")
    chip(s, Inches(4.5), Inches(5.4), "한국어 NLP · 해요체", color=GOLD)
    chip(s, Inches(7.6), Inches(5.4), "AEO + GEO", color=TEAL)
    tx(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
       "출처: 이슈 #1 (SSOT 아키텍처) · 이슈 #4 (한국어 NLP / GEO 기술스택)",
       size=11, color=MUTED)


def s_system_map(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · 시스템 구조", "세 개의 레이어, 하나의 로컬 우선 런타임")
    diagram = (
        "┌─────────────────────────────────────────────────────────────┐\n"
        "│  CLIENT   elbee.yogaman.club · aiegoo.github.io/yoga · curl  │\n"
        "└──────────────────────────────┬──────────────────────────────┘\n"
        "                               │  HTTP / SSE\n"
        "┌──────────────────────────────▼──────────────────────────────┐\n"
        "│  FastAPI :8000   /chat  /chat/stream  /search  /search/...   │\n"
        "│                  /search/locations  /docs  /redoc  /health   │\n"
        "└──────┬─────────────────────────────────────┬────────────────┘\n"
        "       │ RAG path                            │ Location path\n"
        "┌──────▼──────────┐    ┌────────────────────▼───────────────┐\n"
        "│  GeoDataStore   │    │       YogaLocationStore             │\n"
        "│  165 OCR pages  │    │  Haversine · weather · amenity      │\n"
        "│  4 indexes      │    └─────────────────────┬──────────────┘\n"
        "└──────┬──────────┘                          │\n"
        "       │ context chunks               ┌──────▼───────────┐\n"
        "┌──────▼──────────┐                   │   ElbeeAgent      │\n"
        "│  RagService     │                   │   해요체 페르소나   │\n"
        "│  ELBEE_SYSTEM_KO│                   │   조사 자동 결정    │\n"
        "└──────┬──────────┘                   └───────────────────┘\n"
        "┌──────▼──────────────────────────────────────────────────┐\n"
        "│  LLMService → httpx async → Ollama :11434  (llama3)      │\n"
        "└──────────────────────────────────────────────────────────┘"
    )
    code_block(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.7), diagram)
    footer(s, page, total)


def s_korean_nlp(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · 한국어 NLP", "단순 BPE로는 한국어가 깨지는 이유와 해결법")
    items = [
        ("A. 형태소 기반 토크나이징",
         "한국어는 교착어 — 학교에 / 학교도 / 학교만 이 세 토큰이 됩니다.\n"
         "Kiwi · MeCab을 사전 필터로 쓰고, KoBERT · HyperCLOVA-X 계열을 채택했습니다.\n"
         "내부 요가 용어 테스트셋 기준 어간 recall +28%."),
        ("B. 주어 생략 · 문맥 유지",
         "한국어는 주어를 자주 생략합니다. \"그거 맛있더라\" 같은 발화는 turn-1 앵커가 필수.\n"
         "ELBEE_SYSTEM_KO에 해요체 few-shot 예시 + 슬라이딩 히스토리 윈도를 결합.\n"
         "톤앤매너는 표면(chat / invite / search)별로 분리 가능."),
        ("C. 중의성 · 신조어 처리",
         "\"사과하세요\" = apologize 또는 give-an-apple. context window 스캔 +\n"
         "Dynamic Few-shot Learning으로 요청 시점 vector store에서 예시를 끌어옵니다.\n"
         "요린이 · 회복요가 · 인요가 같은 도메인 신조어는 야간 재인덱싱."),
    ]
    y = Inches(2.1)
    for title, body in items:
        card(s, Inches(0.6), y, Inches(12.1), Inches(1.45))
        bullet(s, Inches(0.85), y + Inches(0.18), Inches(11.7), Inches(1.2), title, body)
        y += Inches(1.6)
    footer(s, page, total)


def s_geo_context(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · GEO 컨텍스트 엔진", "위치는 컨텍스트의 끝판왕 — 네 가지 역할")
    items = [
        ("Geofencing",
         "파트너 스튜디오 반경 500m → 실시간 슬롯 / 쿠폰 푸시.\n"
         "/search/locations + Haversine 랭킹 기반."),
        ("Geo-conquesting",
         "경쟁 매장 인근 사용자 감지 → 해요체 톤의 대체 스튜디오 초대\n"
         "메시지를 ElbeeAgent가 즉시 생성."),
        ("상권 분석",
         "유동인구 + 인구통계로 \"어디에 다음 매장을 열 것인가\" 답변.\n"
         "지역 인덱스(성수 · 강남 · 홍대 ...)가 리테일 파일럿에 반영."),
        ("물류 · ETA",
         "실시간 날씨 + 교통으로 실내 / 야외 클래스 재정렬.\n"
         "weather=\"rain\" → weather_indoor=true 자동 필터."),
    ]
    y = Inches(2.1)
    for i, (t, b) in enumerate(items):
        col = i % 2; row = i // 2
        x = Inches(0.6 + col * 6.15); yy = y + Inches(row * 2.3)
        card(s, x, yy, Inches(6.0), Inches(2.1))
        bullet(s, x + Inches(0.25), yy + Inches(0.22), Inches(5.6), Inches(1.7), t, b,
               accent=GOLD if i % 2 else TEAL)
    footer(s, page, total)


def s_ocr_pipeline(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · OCR 인테이크", "도서 스캔에서 28k 단어 지식베이스까지")
    flow = (
        "screenshots/*.png\n"
        "        │\n"
        "        ▼\n"
        "OpenCV 전처리      → grayscale · adaptive threshold · deskew\n"
        "        │\n"
        "        ▼\n"
        "pytesseract (KOR+ENG)  → 페이지별 raw text\n"
        "        │\n"
        "        ▼\n"
        "정제 + 청킹 + 토픽 태깅\n"
        "        │\n"
        "        ▼\n"
        "data/json/generative-engine-optimization/\n"
        "  ├── ocr_database.json     (165 페이지 · 약 28k 단어)\n"
        "  ├── page_index.json        (page_num → 정제 텍스트)\n"
        "  └── keyword_index.json     (역색인)\n"
        "        │\n"
        "        ▼\n"
        "scripts/adapt_content.py  → content/*.md  (Jekyll, frontmatter)\n"
        "scripts/integrate.py       → yoga-chatbot/references/"
    )
    code_block(s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(4.8), flow)
    card(s, Inches(8.4), Inches(2.0), Inches(4.3), Inches(4.8))
    tx(s, Inches(8.65), Inches(2.15), Inches(3.9), Inches(0.4),
       "파이프라인 통계", size=11, bold=True, color=TEAL)
    stats = [
        ("색인 페이지", "165"),
        ("단어 수", "약 28,000"),
        ("인덱스", "page · keyword · platform · region · topic"),
        ("워처", "watchdog 자동 OCR"),
        ("재실행", "단일 명령: python ocr_pipeline.py"),
        ("출력 포맷", "JSON (API) + Markdown (Jekyll)"),
    ]
    yy = Inches(2.6)
    for k, v in stats:
        kv_row(s, Inches(8.65), yy, Inches(4.0), k, v, key_w=Inches(1.4))
        yy += Inches(0.55)
    footer(s, page, total)


def s_stores(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · 인메모리 스토어", "두 개의 스토어, 핫패스 외부 DB 0개")
    card(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8))
    tx(s, Inches(0.85), Inches(2.15), Inches(5.6), Inches(0.4),
       "GeoDataStore", size=16, bold=True, color=TEAL)
    tx(s, Inches(0.85), Inches(2.55), Inches(5.6), Inches(0.4),
       "loader.py · 165 OCR 페이지를 RAM에", size=11, color=MUTED)
    rows = [
        ("소스", "ocr_database.json"),
        ("인덱스", "page · keyword · platform · region · topic"),
        ("조회", "역색인 키워드 → 페이지 리스트 O(1)"),
        ("리프레시", "lifespan startup 훅"),
        ("사용처", "SearchService · RagService"),
    ]
    yy = Inches(3.1)
    for k, v in rows:
        kv_row(s, Inches(0.85), yy, Inches(5.6), k, v, key_w=Inches(1.4))
        yy += Inches(0.55)
    card(s, Inches(6.7), Inches(2.0), Inches(6.0), Inches(4.8))
    tx(s, Inches(6.95), Inches(2.15), Inches(5.6), Inches(0.4),
       "YogaLocationStore", size=16, bold=True, color=GOLD)
    tx(s, Inches(6.95), Inches(2.55), Inches(5.6), Inches(0.4),
       "location_service.py · 서울 8개 시드", size=11, color=MUTED)
    rows = [
        ("소스", "data/yoga_locations.json"),
        ("기하", "Haversine 근접도 (km)"),
        ("필터", "amenity · weather · type · district · tags"),
        ("라우팅", "weather=rain → weather_indoor=true"),
        ("사용처", "/search/locations · ElbeeAgent invite"),
    ]
    yy = Inches(3.1)
    for k, v in rows:
        kv_row(s, Inches(6.95), yy, Inches(5.6), k, v, key_w=Inches(1.4))
        yy += Inches(0.55)
    footer(s, page, total)


def s_rag_llm(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · RAG + LLM", "로컬 우선 생성, OpenAI 의존성 0")
    flow = (
        "ChatRequest{ query, history, geo_coord? }\n"
        "        │\n"
        "        ▼\n"
        "SearchService.search()       → GeoDataStore top-k 페이지\n"
        "        │\n"
        "        ▼\n"
        "RagService.build_context()   → ELBEE_SYSTEM_KO + 페이지 청크\n"
        "        │\n"
        "        ▼\n"
        "LLMService.generate()        → httpx async POST\n"
        "        │                       Ollama :11434  (model = llama3)\n"
        "        ▼\n"
        "ChatResponse{ answer, sources[], brand_logo, language=\"ko\" }"
    )
    code_block(s, Inches(0.6), Inches(2.0), Inches(8.0), Inches(4.0), flow)
    card(s, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.0))
    tx(s, Inches(9.15), Inches(2.15), Inches(3.4), Inches(0.4),
       "왜 로컬인가", size=11, bold=True, color=TEAL)
    notes = [
        "• 호출당 비용 0 — Ollama on-box.",
        "• PII / 헬스 데이터 외부 송신 없음.",
        "• llama3 → HyperCLOVA-X 무중단 교체.",
        "• /chat/stream 으로 SSE 점진 UX.",
    ]
    yy = Inches(2.6)
    for n in notes:
        tx(s, Inches(9.15), yy, Inches(3.4), Inches(0.7), n, size=11, color=INK)
        yy += Inches(0.7)
    footer(s, page, total)


def s_persona(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · 페르소나 엔진", "해요체 + 받침/조사 자동 결정")
    code = (
        "# templates.py — 받침으로 은/는, 이/가, 을/를 자동 선택\n"
        "def josa(word: str, has_batchim: str, no_batchim: str) -> str:\n"
        "    last = word[-1]\n"
        "    code = (ord(last) - 0xAC00) % 28 if 0xAC00 <= ord(last) <= 0xD7A3 else 0\n"
        "    return has_batchim if code else no_batchim\n"
        "\n"
        "# 사용 예\n"
        "msg = f\"{studio}{josa(studio, '은', '는')} 지금 수련 가능해요.\"\n"
        "# → \"성수 elbee 요가 스튜디오는 지금 수련 가능해요.\""
    )
    code_block(s, Inches(0.6), Inches(2.0), Inches(8.0), Inches(2.6), code)
    card(s, Inches(8.9), Inches(2.0), Inches(3.8), Inches(2.6))
    tx(s, Inches(9.15), Inches(2.15), Inches(3.4), Inches(0.4),
       "ELBEE 페르소나", size=11, bold=True, color=GOLD)
    tx(s, Inches(9.15), Inches(2.55), Inches(3.4), Inches(2.0),
       "• 기본 어투 해요체\n• 브랜드 elbee.yogaman.club\n"
       "• 근접 메시지\n• 날씨 인지 라우팅\n• few-shot 톤앤매너",
       size=11, color=INK)
    card(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(1.85))
    tx(s, Inches(0.85), Inches(4.95), Inches(11.6), Inches(0.4),
       "샘플 invite_message (LocationSearchResponse)",
       size=11, bold=True, color=TEAL)
    tx(s, Inches(0.85), Inches(5.4), Inches(11.6), Inches(1.3),
       "안녕하세요, elbee 멤버! 🙏 오늘 성수동은 근처 성수 elbee 요가 스튜디오에서\n"
       "수련하실 기회가 생겼어요. 우천 시 실내 수업으로 자동 매칭됩니다.",
       size=14, color=INK)
    footer(s, page, total)


def s_api_surface(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · API 표면", "11개 라우트 · OpenAPI · Pydantic v2")
    rows = [
        ("POST", "/chat",                "RAG 채팅 — 해요체 답변 + 소스"),
        ("POST", "/chat/stream",         "SSE 스트리밍"),
        ("GET",  "/chat/health",         "LLM + 데이터 스토어 헬스"),
        ("POST", "/search",              "165 페이지 키워드 검색"),
        ("GET",  "/search/suggest",      "프리픽스 자동완성"),
        ("GET",  "/search/filters",      "platform · region · topic 목록"),
        ("GET",  "/search/health",       "인덱스 크기 + 준비 상태"),
        ("POST", "/search/locations",    "근접 / 어메니티 / 날씨 검색"),
        ("GET",  "/docs",                "Swagger UI"),
        ("GET",  "/redoc",               "ReDoc UI"),
        ("GET",  "/health",              "앱 헬스"),
    ]
    yy = Inches(2.0)
    tx(s, Inches(0.6), yy, Inches(0.9), Inches(0.3), "메서드", size=10, bold=True, color=MUTED)
    tx(s, Inches(1.6), yy, Inches(3.5), Inches(0.3), "경로",   size=10, bold=True, color=MUTED)
    tx(s, Inches(5.2), yy, Inches(7.5), Inches(0.3), "설명",   size=10, bold=True, color=MUTED)
    yy += Inches(0.35)
    for m, p, d in rows:
        bar(s, Inches(0.6), yy + Inches(0.05), w=Inches(12.1), h=Inches(0.01), color=DIM)
        accent = TEAL if m == "GET" else GOLD
        tx(s, Inches(0.6), yy + Inches(0.08), Inches(0.9), Inches(0.32),
           m, size=11, bold=True, color=accent)
        tx(s, Inches(1.6), yy + Inches(0.08), Inches(3.5), Inches(0.32),
           p, size=11, color=INK, font="D2Coding")
        tx(s, Inches(5.2), yy + Inches(0.08), Inches(7.5), Inches(0.32),
           d, size=11, color=MUTED)
        yy += Inches(0.4)
    footer(s, page, total)


def s_models(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · 스키마", "Pydantic v2 — 타입 보장 · 검증 · OpenAPI 노출")
    chat = (
        "ChatResponse {\n"
        "  answer: str,\n"
        "  sources: [{ page: int, snippet: str, topic: str }],\n"
        "  model: 'llama3',\n"
        "  brand_logo: HttpUrl,\n"
        "  language: 'ko'\n"
        "}"
    )
    loc = (
        "LocationSearchResponse {\n"
        "  results: [YogaLocationResult{\n"
        "    id, name, type, district, address, lat, lng,\n"
        "    amenities[], tags[], weather_indoor: bool,\n"
        "    rating: float, distance_km: float,\n"
        "    invite_message: str   # 해요체\n"
        "  }],\n"
        "  total, message, brand, brand_logo\n"
        "}"
    )
    code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.7), chat)
    code_block(s, Inches(6.7), Inches(2.0), Inches(6.0), Inches(4.7), loc)
    footer(s, page, total)


def s_twin(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · 크로스플랫폼 트윈", "macOS 개발 → WSL2 / Win11 동등성")
    rows = [
        ("항목",            "macOS",                              "Windows / WSL2"),
        ("Tesseract",       "/opt/homebrew/bin/tesseract",        "C:\\Program Files\\Tesseract-OCR\\tesseract.exe"),
        ("Ollama 엔드포인트", "http://localhost:11434",            "http://host.docker.internal:11434"),
        ("줄바꿈",          "LF",                                  "git config core.autocrlf input"),
        ("Python",          "python3",                             "python  (또는 WSL2의 python3)"),
        ("부트스트랩",       "brew install tesseract ollama",       "wsl --install -d Ubuntu-22.04"),
    ]
    yy = Inches(2.0)
    for i, (a, b, c) in enumerate(rows):
        is_head = (i == 0)
        col_a = MUTED if is_head else TEAL
        col_b = MUTED if is_head else INK
        col_c = MUTED if is_head else INK
        bar(s, Inches(0.6), yy + Inches(0.05), w=Inches(12.1), h=Inches(0.01), color=DIM)
        tx(s, Inches(0.6),  yy + Inches(0.08), Inches(2.4), Inches(0.32),
           a, size=11, bold=True, color=col_a)
        tx(s, Inches(3.0),  yy + Inches(0.08), Inches(4.7), Inches(0.32),
           b, size=11, color=col_b, font="D2Coding" if not is_head else KFONT)
        tx(s, Inches(7.8),  yy + Inches(0.08), Inches(4.9), Inches(0.32),
           c, size=11, color=col_c, font="D2Coding" if not is_head else KFONT)
        yy += Inches(0.55)
    footer(s, page, total)


def s_roadmap(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · 로드맵 · 오픈 이슈", "출시된 것과 다음 것")
    phases = [
        ("Phase 0", "스캐폴드 + OCR 파이프라인",                "완료",  TEAL),
        ("Phase 1", "165 페이지 인테이크 (약 28k 단어 색인)",   "완료",  TEAL),
        ("Phase 2", "elbee.yogaman.club 브랜드 · 8개 스팟",    "완료",  TEAL),
        ("Phase 3", "다중 도서 확장 + 매니페스트",              "다음",  GOLD),
        ("Phase 4", "GitHub Actions CI (OCR 린트 + 스키마)",   "계획",  MUTED),
        ("Phase 5", "고용주용 데모 패키지 + 라이브 배포",        "계획",  MUTED),
    ]
    yy = Inches(2.0)
    for ph, body, tag, col in phases:
        card(s, Inches(0.6), yy, Inches(8.0), Inches(0.65))
        tx(s, Inches(0.85), yy + Inches(0.16), Inches(1.4), Inches(0.4),
           ph, size=13, bold=True, color=col)
        tx(s, Inches(2.3), yy + Inches(0.16), Inches(5.5), Inches(0.4),
           body, size=12, color=INK)
        chip(s, Inches(7.6), yy + Inches(0.18), tag, color=col)
        yy += Inches(0.78)
    card(s, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.7))
    tx(s, Inches(9.15), Inches(2.15), Inches(3.4), Inches(0.4),
       "오픈 이슈", size=11, bold=True, color=GOLD)
    issues = [
        "#1  SSOT 아키텍처 (이 자료의 출처)",
        "#3  60초 세로 스토리보드",
        "#4  한국어 NLP + GEO 기술스택",
    ]
    yy = Inches(2.6)
    for n in issues:
        tx(s, Inches(9.15), yy, Inches(3.4), Inches(0.6),
           "• " + n, size=11, color=INK)
        yy += Inches(0.6)
    footer(s, page, total)


def s_thanks(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               Inches(0.18), H)
    block.line.fill.background(); block.fill.solid()
    block.fill.fore_color.rgb = TEAL
    tx(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.4),
       "로컬 우선으로 만들었습니다.", size=48, bold=True, color=INK)
    tx(s, Inches(0.8), Inches(3.7), Inches(11.5), Inches(1.0),
       "한국어 네이티브. AI가 인용. 확장 준비 완료.",
       size=20, color=MUTED)
    chip(s, Inches(0.8), Inches(5.4), "github.com/aiegoo/aeogeo")
    chip(s, Inches(4.5), Inches(5.4), "elbee.yogaman.club", color=GOLD)
    footer(s, page, total)


def main():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    builders = [
        s_title, s_system_map, s_korean_nlp, s_geo_context, s_ocr_pipeline,
        s_stores, s_rag_llm, s_persona, s_api_surface, s_models, s_twin,
        s_roadmap, s_thanks,
    ]
    total = len(builders)
    s_title(prs)
    for i, fn in enumerate(builders[1:-1], start=2):
        fn(prs, i, total)
    s_thanks(prs, total, total)
    out = Path(__file__).parent / "aeogeo_tech_deepdive_deck_ko.pptx"
    prs.save(out)
    print(f"Wrote {out} {out.stat().st_size} bytes, {total} slides")


if __name__ == "__main__":
    main()
