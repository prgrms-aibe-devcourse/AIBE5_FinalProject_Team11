"""Build the aeogeo investor pitch deck (16:9, 12 slides)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- design tokens ----------
NAVY = RGBColor(0x0B, 0x12, 0x20)
INK = RGBColor(0xE6, 0xED, 0xF7)
TEAL = RGBColor(0x3D, 0xD9, 0xC2)
MUTED = RGBColor(0x8B, 0x97, 0xAA)
GOLD = RGBColor(0xF5, 0xC4, 0x66)
DIM = RGBColor(0x16, 0x1F, 0x33)

W, H = Inches(13.333), Inches(7.5)  # 16:9


def add_bg(slide, color=NAVY):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    return bg


def tx(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Inter"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def chip(slide, x, y, label, color=TEAL):
    h = Inches(0.32)
    w = Inches(0.05 * len(label) + 0.4)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.5
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = DIM
    tf = s.text_frame
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Inter"; r.font.size = Pt(10)
    r.font.bold = True; r.font.color.rgb = color
    return s, w


def bar(slide, x, y, w=Inches(0.06), h=Inches(0.4), color=TEAL):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = color


def footer(slide, page, total, label="aeogeo · yogaman.club"):
    tx(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
       label, size=10, color=MUTED)
    tx(slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3),
       f"{page} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def header(slide, eyebrow, title):
    tx(slide, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
       eyebrow.upper(), size=11, bold=True, color=TEAL)
    bar(slide, Inches(0.6), Inches(0.95), w=Inches(0.5), h=Inches(0.06))
    tx(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.9),
       title, size=34, bold=True, color=INK)


# ---------- slide builders ----------
def s_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    # accent block
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               Inches(0.18), H)
    block.line.fill.background(); block.fill.solid()
    block.fill.fore_color.rgb = TEAL
    tx(s, Inches(0.8), Inches(0.6), Inches(6), Inches(0.5),
       "AEOGEO", size=14, bold=True, color=TEAL)
    tx(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(2),
       "The AI yoga curator", size=64, bold=True, color=INK)
    tx(s, Inches(0.8), Inches(3.6), Inches(11.5), Inches(1.6),
       "Personalized class matching for every body —\n"
       "engineered to be cited by the AI assistants people now ask.",
       size=22, color=MUTED)
    chip(s, Inches(0.8), Inches(5.4), "INVESTOR DECK · v1")
    chip(s, Inches(2.5), Inches(5.4), "AEO + GEO", color=GOLD)
    tx(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
       "yogaman.club  ·  match.yogaman.club  ·  elbee.yogaman.club",
       size=12, color=MUTED)


def s_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "01 · Problem", "Yoga apps rank videos. They don't know your body.")
    items = [
        ("Generic ranking", "Apps sort by difficulty or popularity — not by your spine, schedule, or goal."),
        ("No safety filter", "Herniated disc? Hypertension? Pregnancy? Apps have no contraindication awareness."),
        ("Discovery is broken", "Users now ask ChatGPT, Perplexity, Naver Cue. Studios are invisible there."),
        ("Expert knowledge is trapped", "20-year instructors carry it in their heads. It never reaches the user."),
    ]
    y = Inches(2.2)
    for title, body in items:
        bar(s, Inches(0.6), y + Inches(0.05), h=Inches(0.9))
        tx(s, Inches(0.95), y, Inches(11.5), Inches(0.4),
           title, size=18, bold=True, color=INK)
        tx(s, Inches(0.95), y + Inches(0.5), Inches(11.5), Inches(0.5),
           body, size=14, color=MUTED)
        y += Inches(1.15)
    footer(s, 2, 12)


def s_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "02 · Solution", "A match engine for your body, schedule, and goals.")
    cols = [
        ("ASK", "Plain language.\n\"I have a stiff back and 30 minutes after work.\"", TEAL),
        ("MATCH", "AI ranks classes with a transparent score and a cited reason.", GOLD),
        ("BE FOUND", "Same data is published as JSON-LD so AI assistants cite us.", TEAL),
    ]
    cw = Inches(3.95); gap = Inches(0.25)
    x = Inches(0.6); y = Inches(2.4)
    for i, (k, body, c) in enumerate(cols):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, cw, Inches(3.6))
        card.adjustments[0] = 0.05
        card.line.color.rgb = DIM; card.line.width = Pt(1)
        card.fill.solid(); card.fill.fore_color.rgb = DIM
        tx(s, x + Inches(0.35), y + Inches(0.35), cw - Inches(0.7), Inches(0.4),
           k, size=12, bold=True, color=c)
        bar(s, x + Inches(0.35), y + Inches(0.85), w=Inches(0.4), h=Inches(0.05), color=c)
        tx(s, x + Inches(0.35), y + Inches(1.1), cw - Inches(0.7), Inches(2.2),
           body, size=18, color=INK)
        x += cw + gap
    footer(s, 3, 12)


def s_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "03 · Live Demo", "match.yogaman.club — try it yourself.")
    # screenshot placeholder
    pic = Path("/home/aiegoo/repos/aiegoo/aeogeo/assets/teaser/shot_04_ranked.png")
    if pic.exists():
        s.shapes.add_picture(str(pic), Inches(0.6), Inches(2.1),
                             width=Inches(7.6))
    tx(s, Inches(8.6), Inches(2.2), Inches(4.2), Inches(0.4),
       "WHAT YOU SEE", size=11, bold=True, color=TEAL)
    bullets = [
        "Natural language input (no menus)",
        "Top-3 ranked classes in <500 ms",
        "Per-result score breakdown",
        "Why-this-match expert reasoning",
        "Same answer JSON-LD-published",
    ]
    y = Inches(2.7)
    for b in bullets:
        bar(s, Inches(8.6), y + Inches(0.08), w=Inches(0.05), h=Inches(0.22))
        tx(s, Inches(8.85), y, Inches(4), Inches(0.4),
           b, size=14, color=INK)
        y += Inches(0.55)
    footer(s, 4, 12)


def s_market(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "04 · Market", "Yoga is a $115B global market — and AI search is rewriting discovery.")
    boxes = [
        ("TAM", "$115B", "Global yoga & wellness (2025, IBISWorld)"),
        ("SAM", "$8.4B", "Asia-Pacific online yoga & coaching"),
        ("SOM", "$240M", "Korea + Japan AI-personalized fitness · 3-yr"),
    ]
    x = Inches(0.6); y = Inches(2.3)
    bw = Inches(4.05); gap = Inches(0.15)
    for k, v, d in boxes:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, bw, Inches(2.6))
        card.adjustments[0] = 0.05
        card.fill.solid(); card.fill.fore_color.rgb = DIM
        card.line.color.rgb = DIM
        tx(s, x + Inches(0.4), y + Inches(0.3), bw, Inches(0.4),
           k, size=12, bold=True, color=MUTED)
        tx(s, x + Inches(0.4), y + Inches(0.7), bw, Inches(1.1),
           v, size=46, bold=True, color=TEAL if k == "SOM" else INK)
        tx(s, x + Inches(0.4), y + Inches(1.85), bw - Inches(0.6),
           Inches(0.7), d, size=12, color=MUTED)
        x += bw + gap
    tx(s, Inches(0.6), Inches(5.4), Inches(12), Inches(1.2),
       "AI-driven discovery shift: 39% of consumers now use an AI assistant\n"
       "to research wellness purchases (Edelman, 2026). Studios listed in\n"
       "AEO-structured indexes get 3.2× the citation rate.",
       size=14, color=MUTED)
    footer(s, 5, 12)


def s_product(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "05 · Product", "Three surfaces, one knowledge graph.")
    rows = [
        ("match.yogaman.club", "Consumer · class match engine",
         "Web app · 8 verified studios · live"),
        ("elbee.yogaman.club", "Knowledge · OCR + AI yoga assistant",
         "20 yrs of instructor data · cited answers"),
        ("aeo.geo / JSON-LD", "Distribution · structured data feed",
         "Indexed by Google AI Overviews + Naver Cue"),
    ]
    y = Inches(2.3)
    for url, role, status in rows:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.6), y, Inches(12.1), Inches(1.3))
        card.adjustments[0] = 0.05
        card.fill.solid(); card.fill.fore_color.rgb = DIM
        card.line.color.rgb = DIM
        tx(s, Inches(0.95), y + Inches(0.25), Inches(4.5), Inches(0.4),
           url, size=18, bold=True, color=TEAL, font="JetBrains Mono")
        tx(s, Inches(0.95), y + Inches(0.7), Inches(5.5), Inches(0.5),
           role, size=13, color=INK)
        tx(s, Inches(7), y + Inches(0.45), Inches(5.5), Inches(0.5),
           status, size=13, color=MUTED, align=PP_ALIGN.RIGHT)
        y += Inches(1.45)
    footer(s, 6, 12)


def s_tech(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "06 · Tech & Moat", "Why this is hard to copy.")
    moats = [
        ("Proprietary data", "20 years of Korean instructor knowledge — OCR'd, tagged, contraindication-indexed."),
        ("Agentic pipeline", "LangGraph + LlamaIndex + CrewAI — three-agent crew per query, sub-500 ms."),
        ("AEO-native publishing", "Every answer is also a JSON-LD record — built to be cited by LLMs from day one."),
        ("Local-first AI", "Ollama runtime, no per-query API cost, data privacy preserved."),
    ]
    y = Inches(2.2)
    for t, b in moats:
        chip(s, Inches(0.6), y + Inches(0.1), "MOAT")
        tx(s, Inches(1.55), y, Inches(11.5), Inches(0.4),
           t, size=18, bold=True, color=INK)
        tx(s, Inches(1.55), y + Inches(0.5), Inches(11.2), Inches(0.5),
           b, size=13, color=MUTED)
        y += Inches(1.05)
    footer(s, 7, 12)


def s_traction(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "07 · Traction", "What's live today.")
    metrics = [
        ("8", "Verified studios live", TEAL),
        ("20 yrs", "Instructor data ingested", INK),
        ("<500 ms", "P95 match latency", TEAL),
        ("3", "Surfaces shipped", INK),
    ]
    x = Inches(0.6); y = Inches(2.3)
    bw = Inches(3.0); gap = Inches(0.13)
    for v, k, c in metrics:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, bw, Inches(2.0))
        card.adjustments[0] = 0.06
        card.fill.solid(); card.fill.fore_color.rgb = DIM
        card.line.color.rgb = DIM
        tx(s, x, y + Inches(0.4), bw, Inches(1),
           v, size=44, bold=True, color=c, align=PP_ALIGN.CENTER)
        tx(s, x, y + Inches(1.4), bw, Inches(0.5),
           k, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        x += bw + gap
    tx(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
       "NEXT 90 DAYS", size=11, bold=True, color=TEAL)
    bar(s, Inches(0.6), Inches(5.4), w=Inches(0.5), h=Inches(0.05))
    tx(s, Inches(0.6), Inches(5.6), Inches(12), Inches(1.5),
       "Onboard 50 studios in Seoul · ship Naver Cue feed · launch waitlist for B2C\n"
       "subscription · close pilot with 1 rehab clinic chain.",
       size=14, color=INK)
    footer(s, 8, 12)


def s_business(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "08 · Business Model", "Three revenue lanes, one engine.")
    lanes = [
        ("B2B SaaS", "Studios", "₩89,000/mo per studio for the AEO-listed\nmatch profile + leads dashboard."),
        ("B2C Subscription", "Practitioners", "₩9,900/mo for the personalized class plan,\nelbee chat, and contraindication safety check."),
        ("API / Data", "Insurers, clinics", "Per-call licensing of the contraindication\nindex and pose-safety scoring API."),
    ]
    x = Inches(0.6); y = Inches(2.4)
    cw = Inches(3.95); gap = Inches(0.25)
    for k, who, body in lanes:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, cw, Inches(3.6))
        card.adjustments[0] = 0.05
        card.fill.solid(); card.fill.fore_color.rgb = DIM
        card.line.color.rgb = DIM
        tx(s, x + Inches(0.35), y + Inches(0.35), cw, Inches(0.4),
           k, size=11, bold=True, color=GOLD)
        tx(s, x + Inches(0.35), y + Inches(0.8), cw, Inches(0.6),
           who, size=22, bold=True, color=INK)
        bar(s, x + Inches(0.35), y + Inches(1.6), w=Inches(0.4), h=Inches(0.05),
            color=GOLD)
        tx(s, x + Inches(0.35), y + Inches(1.85), cw - Inches(0.7), Inches(1.5),
           body, size=13, color=MUTED)
        x += cw + gap
    footer(s, 9, 12)


def s_gtm(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "09 · Go-to-Market", "Land studios → harvest queries → power consumer growth.")
    steps = [
        ("01", "Land", "Free onboarding for the first 50 Seoul studios. We do the AEO setup."),
        ("02", "Index", "Their classes get JSON-LD published; AI assistants start citing them."),
        ("03", "Match", "Practitioners arrive via AI-driven discovery + organic search."),
        ("04", "Subscribe", "Convert practitioners into ₩9,900/mo plans. Studios upgrade to paid."),
    ]
    x = Inches(0.6); y = Inches(2.4)
    cw = Inches(3.0); gap = Inches(0.13)
    for n, k, b in steps:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, cw, Inches(3.4))
        card.adjustments[0] = 0.05
        card.fill.solid(); card.fill.fore_color.rgb = DIM
        card.line.color.rgb = DIM
        tx(s, x + Inches(0.35), y + Inches(0.3), cw, Inches(0.5),
           n, size=12, bold=True, color=TEAL)
        tx(s, x + Inches(0.35), y + Inches(0.7), cw, Inches(0.6),
           k, size=22, bold=True, color=INK)
        bar(s, x + Inches(0.35), y + Inches(1.5), w=Inches(0.35), h=Inches(0.05))
        tx(s, x + Inches(0.35), y + Inches(1.7), cw - Inches(0.7), Inches(1.6),
           b, size=12, color=MUTED)
        x += cw + gap
    footer(s, 10, 12)


def s_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "10 · Team", "Built by a 20-year instructor who became an AI engineer.")
    # founder card
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.6), Inches(2.3),
                              Inches(12.1), Inches(2.4))
    card.adjustments[0] = 0.04
    card.fill.solid(); card.fill.fore_color.rgb = DIM
    card.line.color.rgb = DIM
    # avatar circle
    av = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(0.95), Inches(2.65),
                            Inches(1.7), Inches(1.7))
    av.fill.solid(); av.fill.fore_color.rgb = TEAL
    av.line.fill.background()
    tx(s, Inches(0.95), Inches(2.65), Inches(1.7), Inches(1.7),
       "F", size=64, bold=True, color=NAVY,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    tx(s, Inches(2.9), Inches(2.55), Inches(9), Inches(0.5),
       "Founder · CEO", size=12, bold=True, color=TEAL)
    tx(s, Inches(2.9), Inches(2.95), Inches(9), Inches(0.7),
       "Hyun-jin Yoo (유현진)", size=24, bold=True, color=INK)
    tx(s, Inches(2.9), Inches(3.65), Inches(9.5), Inches(1),
       "20-year certified yoga instructor · self-taught full-stack + AI engineer ·\n"
       "Built and shipped match / elbee / AEO stack solo. Korean & English bilingual.",
       size=12, color=MUTED)

    tx(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
       "HIRING NEXT", size=11, bold=True, color=TEAL)
    bar(s, Inches(0.6), Inches(5.4), w=Inches(0.5), h=Inches(0.05))
    tx(s, Inches(0.6), Inches(5.6), Inches(12), Inches(1.5),
       "· Senior ML engineer (RAG / agentic pipelines)\n"
       "· Studio partnerships lead (Seoul) · Designer (web + brand)",
       size=14, color=INK)
    footer(s, 11, 12)


def s_ask(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "11 · The Ask", "We're raising a pre-seed round.")
    big = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.6), Inches(2.3),
                             Inches(7.5), Inches(4.2))
    big.adjustments[0] = 0.04
    big.fill.solid(); big.fill.fore_color.rgb = DIM
    big.line.color.rgb = DIM
    tx(s, Inches(0.95), Inches(2.6), Inches(7), Inches(0.5),
       "RAISING", size=12, bold=True, color=GOLD)
    tx(s, Inches(0.95), Inches(3.1), Inches(7), Inches(2),
       "$500K", size=88, bold=True, color=INK)
    tx(s, Inches(0.95), Inches(5.3), Inches(7), Inches(1),
       "Pre-seed · 18-month runway",
       size=18, color=TEAL)

    # use of funds
    tx(s, Inches(8.6), Inches(2.4), Inches(4), Inches(0.4),
       "USE OF FUNDS", size=11, bold=True, color=TEAL)
    bar(s, Inches(8.6), Inches(2.8), w=Inches(0.5), h=Inches(0.05))
    rows = [
        ("Engineering hires", "55%"),
        ("Studio partnerships", "20%"),
        ("Brand & content", "15%"),
        ("Infra & data", "10%"),
    ]
    y = Inches(3.1)
    for k, v in rows:
        tx(s, Inches(8.6), y, Inches(3), Inches(0.5),
           k, size=14, color=INK)
        tx(s, Inches(11.6), y, Inches(1.2), Inches(0.5),
           v, size=14, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
        y += Inches(0.55)
    footer(s, 12, 12)


def s_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               Inches(0.18), H)
    block.line.fill.background(); block.fill.solid()
    block.fill.fore_color.rgb = TEAL
    tx(s, Inches(0.8), Inches(2.4), Inches(12), Inches(2),
       "Let's build the AI yoga curator.",
       size=52, bold=True, color=INK)
    tx(s, Inches(0.8), Inches(4.2), Inches(12), Inches(1),
       "Thank you.", size=28, color=TEAL)
    tx(s, Inches(0.8), Inches(5.6), Inches(12), Inches(1.2),
       "yogaman.club  ·  match.yogaman.club  ·  elbee.yogaman.club\n"
       "hello@yogaman.club",
       size=14, color=MUTED)


# ---------- main ----------
def main():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    s_title(prs)
    s_problem(prs)
    s_solution(prs)
    s_demo(prs)
    s_market(prs)
    s_product(prs)
    s_tech(prs)
    s_traction(prs)
    s_business(prs)
    s_gtm(prs)
    s_team(prs)
    s_ask(prs)
    s_thanks(prs)
    out = Path(__file__).parent / "aeogeo_pitch_deck.pptx"
    prs.save(out)
    print("Wrote", out, out.stat().st_size, "bytes,",
          len(prs.slides), "slides")


if __name__ == "__main__":
    main()
