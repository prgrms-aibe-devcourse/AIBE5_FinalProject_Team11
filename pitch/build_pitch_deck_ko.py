"""Build the aeogeo investor pitch deck — Korean (16:9, 13 slides)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x0B, 0x12, 0x20)
INK = RGBColor(0xE6, 0xED, 0xF7)
TEAL = RGBColor(0x3D, 0xD9, 0xC2)
MUTED = RGBColor(0x8B, 0x97, 0xAA)
GOLD = RGBColor(0xF5, 0xC4, 0x66)
DIM = RGBColor(0x16, 0x1F, 0x33)

W, H = Inches(13.333), Inches(7.5)
KFONT = "Pretendard"  # falls back to system Korean font if absent


def add_bg(slide, color=NAVY):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    return bg


def tx(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=KFONT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def chip(slide, x, y, label, color=TEAL):
    h = Inches(0.32)
    w = Inches(0.08 * len(label) + 0.5)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.5
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = DIM
    tf = s.text_frame
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = KFONT; r.font.size = Pt(10)
    r.font.bold = True; r.font.color.rgb = color
    return s, w


def bar(slide, x, y, w=Inches(0.06), h=Inches(0.4), color=TEAL):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = color


def footer(slide, page, total, label="aeogeo · yogaman.club"):
    tx(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
       label, size=10, color=MUTED)
    tx(slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3),
       f"{page} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def header(slide, eyebrow, title):
    tx(slide, Inches(0.6), Inches(0.5), Inches(10), Inches(0.4),
       eyebrow, size=11, bold=True, color=TEAL)
    bar(slide, Inches(0.6), Inches(0.95), w=Inches(0.5), h=Inches(0.06))
    tx(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.9),
       title, size=32, bold=True, color=INK)


# ---------- slides ----------
def s_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               Inches(0.18), H)
    block.line.fill.background(); block.fill.solid()
    block.fill.fore_color.rgb = TEAL
    tx(s, Inches(0.8), Inches(0.6), Inches(6), Inches(0.5),
       "AEOGEO", size=14, bold=True, color=TEAL)
    tx(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(2),
       "AI 요가 큐레이터", size=64, bold=True, color=INK)
    tx(s, Inches(0.8), Inches(3.6), Inches(11.5), Inches(1.8),
       "당신의 몸·일정·목표에 맞춘 요가 매칭.\n"
       "그리고 AI 어시스턴트가 추천하도록 설계된 데이터.",
       size=22, color=MUTED)
    chip(s, Inches(0.8), Inches(5.6), "투자제안서 · v1")
    chip(s, Inches(2.8), Inches(5.6), "AEO + GEO", color=GOLD)
    tx(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
       "yogaman.club  ·  match.yogaman.club  ·  elbee.yogaman.club",
       size=12, color=MUTED)


def s_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "01 · 문제", "요가 앱은 영상을 줄세울 뿐, 사용자의 몸을 모릅니다.")
    items = [
        ("획일적 추천", "난이도와 인기 순으로 정렬할 뿐, 척추·일정·목표를 보지 않습니다."),
        ("안전 필터 부재", "디스크, 고혈압, 임신 등 금기사항을 인지하지 못합니다."),
        ("발견 채널의 변화", "사용자는 이제 ChatGPT·Perplexity·Naver Cue에 묻습니다.\n그 안에서 스튜디오는 보이지 않습니다."),
        ("전문가 지식의 단절", "20년차 강사의 노하우가 사용자에게 닿지 않습니다."),
    ]
    y = Inches(2.2)
    for title, body in items:
        bar(s, Inches(0.6), y + Inches(0.05), h=Inches(0.95))
        tx(s, Inches(0.95), y, Inches(11.5), Inches(0.4),
           title, size=18, bold=True, color=INK)
        tx(s, Inches(0.95), y + Inches(0.5), Inches(11.5), Inches(0.6),
           body, size=13, color=MUTED)
        y += Inches(1.15)
    footer(s, 2, 13)


def s_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "02 · 솔루션", "몸·일정·목표를 위한 매칭 엔진.")
    cols = [
        ("ASK · 질문", "자연어로 입력하세요.\n\"허리가 뻐근한데 퇴근 후 30분 가능해요.\"", TEAL),
        ("MATCH · 추천", "AI가 점수와 근거를 함께 보여주는\n맞춤 클래스를 추천합니다.", GOLD),
        ("BE FOUND · 노출", "동일한 데이터를 JSON-LD로 발행해\nAI 어시스턴트가 우리를 인용합니다.", TEAL),
    ]
    cw = Inches(3.95); gap = Inches(0.25)
    x = Inches(0.6); y = Inches(2.4)
    for k, body, c in cols:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, cw, Inches(3.6))
        card.adjustments[0] = 0.05
        card.line.color.rgb = DIM; card.line.width = Pt(1)
        card.fill.solid(); card.fill.fore_color.rgb = DIM
        tx(s, x + Inches(0.35), y + Inches(0.35), cw - Inches(0.7), Inches(0.4),
           k, size=12, bold=True, color=c)
        bar(s, x + Inches(0.35), y + Inches(0.85), w=Inches(0.4), h=Inches(0.05), color=c)
        tx(s, x + Inches(0.35), y + Inches(1.1), cw - Inches(0.7), Inches(2.2),
           body, size=16, color=INK)
        x += cw + gap
    footer(s, 3, 13)


def s_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "03 · 라이브 데모", "match.yogaman.club — 지금 바로 사용해 보세요.")

    BASE = Path("/home/aiegoo/repos/aiegoo/aeogeo/assets/teaser")
    shots = [
        (BASE / "shot_03_typing.png",    "① 목표·조건 입력"),
        (BASE / "shot_04_ranked.png",    "② Top-3 추천"),
        (BASE / "shot_05_breakdown.png", "③ 점수 분해"),
        (BASE / "shot_07_jsonld.png",    "④ JSON-LD 발행"),
    ]
    iw = Inches(3.75)
    ih = Inches(3.75 * 9 / 16)   # 2.11" — 16:9 aspect
    xs = [Inches(0.6), Inches(4.45)]
    ys = [Inches(2.1), Inches(4.5)]
    for idx, (pic, caption) in enumerate(shots):
        cx, cy = xs[idx % 2], ys[idx // 2]
        if pic.exists():
            s.shapes.add_picture(str(pic), cx, cy, width=iw)
        else:
            rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, iw, ih)
            rect.fill.solid(); rect.fill.fore_color.rgb = DIM
            rect.line.fill.background()
        tx(s, cx, cy + ih + Inches(0.04), iw, Inches(0.22),
           caption, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    tx(s, Inches(8.6), Inches(2.2), Inches(4.2), Inches(0.4),
       "보이는 것", size=11, bold=True, color=TEAL)
    bullets = [
        "자연어 입력 (메뉴 없음)",
        "500ms 이내 Top-3 추천",
        "결과별 점수 분해",
        "왜 이 클래스인지 전문가 근거",
        "동일 답변을 JSON-LD로 발행",
    ]
    y = Inches(2.7)
    for b in bullets:
        bar(s, Inches(8.6), y + Inches(0.08), w=Inches(0.05), h=Inches(0.22))
        tx(s, Inches(8.85), y, Inches(4), Inches(0.4),
           b, size=14, color=INK)
        y += Inches(0.55)
    footer(s, 4, 13)


def s_market(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "04 · 시장", "글로벌 요가 시장 1,150억 달러. AI 검색이 발견의 룰을 다시 씁니다.")
    boxes = [
        ("TAM", "$115B", "글로벌 요가·웰니스 (2025, IBISWorld)"),
        ("SAM", "$8.4B", "아시아·태평양 온라인 요가·코칭"),
        ("SOM", "$240M", "한·일 AI 맞춤 피트니스 · 3년"),
    ]
    x = Inches(0.6); y = Inches(2.3)
    bw = Inches(4.05); gap = Inches(0.15)
    for k, v, d in boxes:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, bw, Inches(2.6))
        card.adjustments[0] = 0.05
        card.fill.solid(); card.fill.fore_color.rgb = DIM
        card.line.color.rgb = DIM
        tx(s, x + Inches(0.4), y + Inches(0.3), bw, Inches(0.4),
           k, size=12, bold=True, color=MUTED)
        tx(s, x + Inches(0.4), y + Inches(0.7), bw, Inches(1.1),
           v, size=46, bold=True, color=TEAL if k == "SOM" else INK)
        tx(s, x + Inches(0.4), y + Inches(1.85), bw - Inches(0.6),
           Inches(0.7), d, size=12, color=MUTED)
        x += bw + gap
    tx(s, Inches(0.6), Inches(5.4), Inches(12), Inches(1.4),
       "AI 기반 발견의 전환: 소비자의 39%가 웰니스 구매 전\n"
       "AI 어시스턴트로 정보를 탐색합니다 (Edelman, 2026).\n"
       "AEO 구조화된 인덱스에 등록된 스튜디오의 인용률은 3.2배 높습니다.",
       size=13, color=MUTED)
    footer(s, 5, 13)


def s_product(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "05 · 제품", "세 개의 표면, 하나의 지식 그래프.")
    rows = [
        ("match.yogaman.club", "소비자 · 클래스 매칭 엔진",
         "웹앱 · 검증된 8개 스튜디오 · 운영 중"),
        ("elbee.yogaman.club", "지식 · OCR + AI 요가 어시스턴트",
         "20년 강사 데이터 · 인용 기반 답변"),
        ("aeo.geo / JSON-LD", "유통 · 구조화 데이터 피드",
         "Google AI Overviews + Naver Cue 색인"),
    ]
    y = Inches(2.3)
    for url, role, status in rows:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.6), y, Inches(12.1), Inches(1.3))
        card.adjustments[0] = 0.05
        card.fill.solid(); card.fill.fore_color.rgb = DIM
        card.line.color.rgb = DIM
        tx(s, Inches(0.95), y + Inches(0.25), Inches(4.5), Inches(0.4),
           url, size=18, bold=True, color=TEAL, font="JetBrains Mono")
        tx(s, Inches(0.95), y + Inches(0.7), Inches(5.5), Inches(0.5),
           role, size=13, color=INK)
        tx(s, Inches(7), y + Inches(0.45), Inches(5.5), Inches(0.5),
           status, size=13, color=MUTED, align=PP_ALIGN.RIGHT)
        y += Inches(1.45)
    footer(s, 6, 13)


def s_tech(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "06 · 기술과 해자", "복제하기 어려운 이유.")
    moats = [
        ("독자 데이터", "20년치 한국 강사 노하우를 OCR·태깅·금기사항 인덱싱."),
        ("에이전틱 파이프라인", "LangGraph + LlamaIndex + CrewAI — 쿼리당 3-에이전트, 500ms 이하."),
        ("AEO 네이티브 발행", "모든 답변이 JSON-LD 레코드. 처음부터 LLM이 인용하도록 설계."),
        ("로컬 우선 AI", "Ollama 런타임 — 쿼리당 API 비용 0, 데이터 프라이버시 보존."),
    ]
    y = Inches(2.2)
    for t, b in moats:
        chip(s, Inches(0.6), y + Inches(0.1), "MOAT")
        tx(s, Inches(1.7), y, Inches(11.5), Inches(0.4),
           t, size=18, bold=True, color=INK)
        tx(s, Inches(1.7), y + Inches(0.5), Inches(11), Inches(0.5),
           b, size=13, color=MUTED)
        y += Inches(1.05)
    footer(s, 7, 13)


def s_traction(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "07 · 트랙션", "지금 운영 중인 것.")
    metrics = [
        ("8", "검증 스튜디오", TEAL),
        ("20년", "강사 데이터", INK),
        ("<500ms", "P95 매칭 지연", TEAL),
        ("3", "출시 채널", INK),
    ]
    x = Inches(0.6); y = Inches(2.3)
    bw = Inches(3.0); gap = Inches(0.13)
    for v, k, c in metrics:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, bw, Inches(2.0))
        card.adjustments[0] = 0.06
        card.fill.solid(); card.fill.fore_color.rgb = DIM
        card.line.color.rgb = DIM
        tx(s, x, y + Inches(0.4), bw, Inches(1),
           v, size=42, bold=True, color=c, align=PP_ALIGN.CENTER)
        tx(s, x, y + Inches(1.4), bw, Inches(0.5),
           k, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        x += bw + gap
    tx(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
       "향후 90일", size=11, bold=True, color=TEAL)
    bar(s, Inches(0.6), Inches(5.4), w=Inches(0.5), h=Inches(0.05))
    tx(s, Inches(0.6), Inches(5.6), Inches(12), Inches(1.5),
       "서울 50개 스튜디오 온보딩 · Naver Cue 피드 출시 · B2C 구독 대기열\n"
       "오픈 · 재활클리닉 체인 1곳 파일럿 체결.",
       size=14, color=INK)
    footer(s, 8, 13)


def s_business(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "08 · 비즈니스 모델", "세 개의 매출 라인, 하나의 엔진.")
    lanes = [
        ("B2B SaaS", "스튜디오", "월 89,000원 / 스튜디오\nAEO 매칭 프로필 + 리드 대시보드"),
        ("B2C 구독", "수련자", "월 9,900원\n맞춤 클래스 플랜 + elbee 챗 + 안전 필터"),
        ("API · 데이터", "보험·재활클리닉", "금기사항 인덱스 + 자세 안전성\n스코어링 API 호출당 과금"),
    ]
    x = Inches(0.6); y = Inches(2.4)
    cw = Inches(3.95); gap = Inches(0.25)
    for k, who, body in lanes:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, cw, Inches(3.6))
        card.adjustments[0] = 0.05
        card.fill.solid(); card.fill.fore_color.rgb = DIM
        card.line.color.rgb = DIM
        tx(s, x + Inches(0.35), y + Inches(0.35), cw, Inches(0.4),
           k, size=11, bold=True, color=GOLD)
        tx(s, x + Inches(0.35), y + Inches(0.8), cw, Inches(0.6),
           who, size=22, bold=True, color=INK)
        bar(s, x + Inches(0.35), y + Inches(1.6), w=Inches(0.4), h=Inches(0.05),
            color=GOLD)
        tx(s, x + Inches(0.35), y + Inches(1.85), cw - Inches(0.7), Inches(1.5),
           body, size=13, color=MUTED)
        x += cw + gap
    footer(s, 9, 13)


def s_gtm(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "09 · 시장 진입 전략", "스튜디오 확보 → 쿼리 수확 → 소비자 성장.")
    steps = [
        ("01", "확보", "서울 첫 50개 스튜디오 무료 온보딩.\nAEO 셋업을 우리가 직접 수행."),
        ("02", "색인", "JSON-LD로 발행되어 AI 어시스턴트가\n클래스를 인용하기 시작."),
        ("03", "매칭", "AI 발견 + 자연 검색을 통해\n수련자가 유입됩니다."),
        ("04", "구독", "월 9,900원으로 전환.\n스튜디오는 유료 플랜으로 업그레이드."),
    ]
    x = Inches(0.6); y = Inches(2.4)
    cw = Inches(3.0); gap = Inches(0.13)
    for n, k, b in steps:
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, cw, Inches(3.4))
        card.adjustments[0] = 0.05
        card.fill.solid(); card.fill.fore_color.rgb = DIM
        card.line.color.rgb = DIM
        tx(s, x + Inches(0.35), y + Inches(0.3), cw, Inches(0.5),
           n, size=12, bold=True, color=TEAL)
        tx(s, x + Inches(0.35), y + Inches(0.7), cw, Inches(0.6),
           k, size=22, bold=True, color=INK)
        bar(s, x + Inches(0.35), y + Inches(1.5), w=Inches(0.35), h=Inches(0.05))
        tx(s, x + Inches(0.35), y + Inches(1.7), cw - Inches(0.7), Inches(1.6),
           b, size=12, color=MUTED)
        x += cw + gap
    footer(s, 10, 13)


def s_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "10 · 팀", "20년차 요가 강사가 직접 만든 AI 엔진.")
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.6), Inches(2.3),
                              Inches(12.1), Inches(2.4))
    card.adjustments[0] = 0.04
    card.fill.solid(); card.fill.fore_color.rgb = DIM
    card.line.color.rgb = DIM
    av = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(0.95), Inches(2.65),
                            Inches(1.7), Inches(1.7))
    av.fill.solid(); av.fill.fore_color.rgb = TEAL
    av.line.fill.background()
    tx(s, Inches(0.95), Inches(2.65), Inches(1.7), Inches(1.7),
       "F", size=64, bold=True, color=NAVY,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    tx(s, Inches(2.9), Inches(2.55), Inches(9), Inches(0.5),
       "창업자 · CEO", size=12, bold=True, color=TEAL)
    tx(s, Inches(2.9), Inches(2.95), Inches(9), Inches(0.7),
       "유현진 (Hyun-jin Yoo)", size=24, bold=True, color=INK)
    tx(s, Inches(2.9), Inches(3.65), Inches(9.5), Inches(1),
       "20년 경력 공인 요가 강사 · 독학 풀스택·AI 엔지니어\n"
       "match·elbee·AEO 스택을 단독으로 설계·출시. 한·영 이중언어.",
       size=12, color=MUTED)

    tx(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
       "다음 채용", size=11, bold=True, color=TEAL)
    bar(s, Inches(0.6), Inches(5.4), w=Inches(0.5), h=Inches(0.05))
    tx(s, Inches(0.6), Inches(5.6), Inches(12), Inches(1.5),
       "· 시니어 ML 엔지니어 (RAG / 에이전틱 파이프라인)\n"
       "· 스튜디오 파트너십 리드 (서울) · 디자이너 (웹 + 브랜드)",
       size=14, color=INK)
    footer(s, 11, 13)


def s_ask(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, "11 · 투자 요청", "프리시드 라운드 모집 중.")
    big = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.6), Inches(2.3),
                             Inches(7.5), Inches(4.2))
    big.adjustments[0] = 0.04
    big.fill.solid(); big.fill.fore_color.rgb = DIM
    big.line.color.rgb = DIM
    tx(s, Inches(0.95), Inches(2.6), Inches(7), Inches(0.5),
       "모집 금액", size=12, bold=True, color=GOLD)
    tx(s, Inches(0.95), Inches(3.1), Inches(7), Inches(2),
       "$500K", size=88, bold=True, color=INK)
    tx(s, Inches(0.95), Inches(5.3), Inches(7), Inches(1),
       "프리시드 · 18개월 런웨이",
       size=18, color=TEAL)

    tx(s, Inches(8.6), Inches(2.4), Inches(4), Inches(0.4),
       "사용처", size=11, bold=True, color=TEAL)
    bar(s, Inches(8.6), Inches(2.8), w=Inches(0.5), h=Inches(0.05))
    rows = [
        ("엔지니어링 채용", "55%"),
        ("스튜디오 파트너십", "20%"),
        ("브랜드·콘텐츠", "15%"),
        ("인프라·데이터", "10%"),
    ]
    y = Inches(3.1)
    for k, v in rows:
        tx(s, Inches(8.6), y, Inches(3), Inches(0.5),
           k, size=14, color=INK)
        tx(s, Inches(11.6), y, Inches(1.2), Inches(0.5),
           v, size=14, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
        y += Inches(0.55)
    footer(s, 12, 13)


def s_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               Inches(0.18), H)
    block.line.fill.background(); block.fill.solid()
    block.fill.fore_color.rgb = TEAL
    tx(s, Inches(0.8), Inches(2.4), Inches(12), Inches(2),
       "함께 AI 요가 큐레이터를 만들어요.",
       size=46, bold=True, color=INK)
    tx(s, Inches(0.8), Inches(4.2), Inches(12), Inches(1),
       "감사합니다.", size=28, color=TEAL)
    tx(s, Inches(0.8), Inches(5.6), Inches(12), Inches(1.2),
       "yogaman.club  ·  match.yogaman.club  ·  elbee.yogaman.club\n"
       "hello@yogaman.club",
       size=14, color=MUTED)


def main():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    s_title(prs); s_problem(prs); s_solution(prs); s_demo(prs)
    s_market(prs); s_product(prs); s_tech(prs); s_traction(prs)
    s_business(prs); s_gtm(prs); s_team(prs); s_ask(prs); s_thanks(prs)
    out = Path(__file__).parent / "aeogeo_pitch_deck_ko.pptx"
    prs.save(out)
    print("Wrote", out, out.stat().st_size, "bytes,",
          len(prs.slides), "slides")


if __name__ == "__main__":
    main()
