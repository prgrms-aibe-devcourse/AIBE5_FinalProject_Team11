"""Detailed engineering deep-dive deck (EN). Sourced from GitHub issues #1, #4
and the SSOT architecture in the aiegoo/aeogeo repo. 16:9, 12 slides."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from build_pitch_deck import (
    NAVY, INK, TEAL, MUTED, GOLD, DIM,
    W, H, add_bg, tx, chip, bar, footer, header,
)

FONT = "Inter"


# ---------- shared layout helpers ----------
def card(slide, x, y, w, h, fill=DIM):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.05
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s


def kv_row(slide, x, y, w, key, val, key_w=Inches(2.2)):
    tx(slide, x, y, key_w, Inches(0.32), key, size=11, bold=True, color=TEAL, font=FONT)
    tx(slide, x + key_w, y, w - key_w, Inches(0.32), val, size=11, color=INK, font=FONT)


def bullet(slide, x, y, w, h, title, body, accent=TEAL):
    bar(slide, x, y + Inches(0.04), w=Inches(0.06), h=Inches(0.42), color=accent)
    tx(slide, x + Inches(0.18), y, w - Inches(0.2), Inches(0.4),
       title, size=14, bold=True, color=INK, font=FONT)
    tx(slide, x + Inches(0.18), y + Inches(0.45), w - Inches(0.2), h - Inches(0.45),
       body, size=11, color=MUTED, font=FONT)


def code_block(slide, x, y, w, h, text):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = DIM
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = tf.margin_bottom = Inches(0.14)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = line
        r.font.name = "Consolas"; r.font.size = Pt(10)
        r.font.color.rgb = TEAL if line.startswith(("→", "│", "├", "└", "┌", "┐")) else INK


# ---------- slide builders ----------
def s_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               Inches(0.18), H)
    block.line.fill.background(); block.fill.solid()
    block.fill.fore_color.rgb = TEAL
    tx(s, Inches(0.8), Inches(0.6), Inches(6), Inches(0.5),
       "AEOGEO · ENGINEERING", size=14, bold=True, color=TEAL, font=FONT)
    tx(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.4),
       "Engineering deep-dive", size=58, bold=True, color=INK, font=FONT)
    tx(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(1.6),
       "How we turn a 165-page domain book into a Korean-native\n"
       "RAG API that AI assistants actually cite.", size=22, color=MUTED, font=FONT)
    chip(s, Inches(0.8), Inches(5.4), "FastAPI · Ollama · Tesseract")
    chip(s, Inches(4.0), Inches(5.4), "Korean NLP · 해요체", color=GOLD)
    chip(s, Inches(6.7), Inches(5.4), "AEO + GEO", color=TEAL)
    tx(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
       "Sourced from issues #1 (SSOT architecture) and #4 (Korean NLP / GEO techstack)",
       size=11, color=MUTED, font=FONT)


def s_system_map(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · System map", "Three layers, one local-first runtime")
    diagram = (
        "┌─────────────────────────────────────────────────────────────┐\n"
        "│  CLIENT   elbee.yogaman.club · aiegoo.github.io/yoga · curl  │\n"
        "└──────────────────────────────┬──────────────────────────────┘\n"
        "                               │  HTTP / SSE\n"
        "┌──────────────────────────────▼──────────────────────────────┐\n"
        "│  FastAPI :8000   /chat  /chat/stream  /search  /search/...   │\n"
        "│                  /search/locations  /docs  /redoc  /health   │\n"
        "└──────┬─────────────────────────────────────┬────────────────┘\n"
        "       │ RAG path                            │ Location path\n"
        "┌──────▼──────────┐    ┌────────────────────▼───────────────┐\n"
        "│  GeoDataStore   │    │       YogaLocationStore             │\n"
        "│  165 OCR pages  │    │  Haversine · weather · amenity      │\n"
        "│  4 indexes      │    └─────────────────────┬──────────────┘\n"
        "└──────┬──────────┘                          │\n"
        "       │ context chunks               ┌──────▼───────────┐\n"
        "┌──────▼──────────┐                   │   ElbeeAgent      │\n"
        "│  RagService     │                   │   해요체 persona   │\n"
        "│  ELBEE_SYSTEM_KO│                   │   조사 resolver    │\n"
        "└──────┬──────────┘                   └───────────────────┘\n"
        "┌──────▼──────────────────────────────────────────────────┐\n"
        "│  LLMService → httpx async → Ollama :11434  (llama3)      │\n"
        "└──────────────────────────────────────────────────────────┘"
    )
    code_block(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.7), diagram)
    footer(s, page, total)


def s_korean_nlp(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · Korean NLP", "Why naive BPE breaks 한국어 — and how we fix it")
    items = [
        ("A. Morpheme-aware tokenization",
         "Korean is agglutinative: 학교에 / 학교도 / 학교만 are three tokens for one stem.\n"
         "Pre-filter with Kiwi or MeCab; ground generation in KoBERT / HyperCLOVA-X aware models.\n"
         "Result: stem recall + 28% on internal yoga-term test set."),
        ("B. Subject-drop & context retention",
         "한국어 drops subjects. \"그거 맛있더라\" requires turn-1 anchoring.\n"
         "Few-shot 해요체 exemplars in ELBEE_SYSTEM_KO + sliding history window.\n"
         "Tone-and-manner profile selectable per surface (chat vs. invite vs. search)."),
        ("C. Polysemy & neologism handling",
         "\"사과하세요\" = apologize OR give-an-apple. Disambiguated by context-window scan +\n"
         "Dynamic Few-shot Learning that pulls fresh exemplars from the vector store at request time.\n"
         "Yoga-domain neologisms (요린이, 회복요가, 인요가) re-indexed nightly."),
    ]
    y = Inches(2.1)
    for title, body in items:
        card(s, Inches(0.6), y, Inches(12.1), Inches(1.45))
        bullet(s, Inches(0.85), y + Inches(0.18), Inches(11.7), Inches(1.2), title, body)
        y += Inches(1.6)
    footer(s, page, total)


def s_geo_context(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · GEO context engine", "Location is the ultimate context — four jobs")
    items = [
        ("Geofencing",
         "500 m radius around a partner studio → push a real-time slot or coupon.\n"
         "Backed by /search/locations with Haversine ranking."),
        ("Geo-conquesting",
         "User detected near a competitor → trigger a comparable studio invite\n"
         "with a tone-matched 해요체 message from ElbeeAgent."),
        ("Site selection",
         "Foot-traffic + demographic overlay answers \"where to open next.\"\n"
         "District index (성수, 강남, 홍대...) feeds retail pilots."),
        ("Logistics & ETA",
         "Live weather + traffic re-ranks indoor vs. outdoor classes.\n"
         "weather=\"rain\" filters to weather_indoor=true only."),
    ]
    y = Inches(2.1)
    for i, (t, b) in enumerate(items):
        col = i % 2; row = i // 2
        x = Inches(0.6 + col * 6.15); yy = y + Inches(row * 2.3)
        card(s, x, yy, Inches(6.0), Inches(2.1))
        bullet(s, x + Inches(0.25), yy + Inches(0.22), Inches(5.6), Inches(1.7), t, b,
               accent=GOLD if i % 2 else TEAL)
    footer(s, page, total)


def s_ocr_pipeline(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · OCR intake", "From book scans to a 28k-word knowledge base")
    flow = (
        "screenshots/*.png\n"
        "        │\n"
        "        ▼\n"
        "OpenCV preprocess   → grayscale · adaptive threshold · deskew\n"
        "        │\n"
        "        ▼\n"
        "pytesseract (KOR+ENG)  → raw text per page\n"
        "        │\n"
        "        ▼\n"
        "clean + chunk + topic-tag\n"
        "        │\n"
        "        ▼\n"
        "data/json/generative-engine-optimization/\n"
        "  ├── ocr_database.json     (165 pages · ~28k words)\n"
        "  ├── page_index.json        (page_num → cleaned text)\n"
        "  └── keyword_index.json     (inverted index)\n"
        "        │\n"
        "        ▼\n"
        "scripts/adapt_content.py  → content/*.md  (Jekyll, frontmatter)\n"
        "scripts/integrate.py       → yoga-chatbot/references/"
    )
    code_block(s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(4.8), flow)
    # right column: stats
    card(s, Inches(8.4), Inches(2.0), Inches(4.3), Inches(4.8))
    tx(s, Inches(8.65), Inches(2.15), Inches(3.9), Inches(0.4),
       "PIPELINE STATS", size=11, bold=True, color=TEAL, font=FONT)
    stats = [
        ("Pages indexed", "165"),
        ("Words", "≈ 28,000"),
        ("Indexes built", "page · keyword · platform · region · topic"),
        ("Watcher", "watchdog auto-OCR on image drop"),
        ("Re-run cost", "single command: python ocr_pipeline.py"),
        ("Output formats", "JSON (API) + Markdown (Jekyll)"),
    ]
    yy = Inches(2.6)
    for k, v in stats:
        kv_row(s, Inches(8.65), yy, Inches(4.0), k, v, key_w=Inches(1.7))
        yy += Inches(0.55)
    footer(s, page, total)


def s_stores(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · In-memory stores", "Two stores, zero external DB on the hot path")
    # left card
    card(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8))
    tx(s, Inches(0.85), Inches(2.15), Inches(5.6), Inches(0.4),
       "GeoDataStore", size=16, bold=True, color=TEAL, font=FONT)
    tx(s, Inches(0.85), Inches(2.55), Inches(5.6), Inches(0.4),
       "loader.py · 165 OCR pages in RAM", size=11, color=MUTED, font=FONT)
    rows = [
        ("Source", "ocr_database.json"),
        ("Indexes", "page · keyword · platform · region · topic"),
        ("Lookup", "O(1) keyword → page list (inverted)"),
        ("Refresh", "lifespan startup hook"),
        ("Used by", "SearchService · RagService"),
    ]
    yy = Inches(3.1)
    for k, v in rows:
        kv_row(s, Inches(0.85), yy, Inches(5.6), k, v, key_w=Inches(1.4))
        yy += Inches(0.55)
    # right card
    card(s, Inches(6.7), Inches(2.0), Inches(6.0), Inches(4.8))
    tx(s, Inches(6.95), Inches(2.15), Inches(5.6), Inches(0.4),
       "YogaLocationStore", size=16, bold=True, color=GOLD, font=FONT)
    tx(s, Inches(6.95), Inches(2.55), Inches(5.6), Inches(0.4),
       "location_service.py · 8 Seoul spots seed", size=11, color=MUTED, font=FONT)
    rows = [
        ("Source", "data/yoga_locations.json"),
        ("Geometry", "Haversine proximity (km)"),
        ("Filters", "amenity · weather · type · district · tags"),
        ("Routing", "weather=rain → weather_indoor=true only"),
        ("Used by", "/search/locations · ElbeeAgent invite"),
    ]
    yy = Inches(3.1)
    for k, v in rows:
        kv_row(s, Inches(6.95), yy, Inches(5.6), k, v, key_w=Inches(1.4))
        yy += Inches(0.55)
    footer(s, page, total)


def s_rag_llm(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · RAG + LLM", "Local-first generation, no OpenAI dependency")
    # diagram
    flow = (
        "ChatRequest{ query, history, geo_coord? }\n"
        "        │\n"
        "        ▼\n"
        "SearchService.search()       → top-k pages from GeoDataStore\n"
        "        │\n"
        "        ▼\n"
        "RagService.build_context()   → ELBEE_SYSTEM_KO + page chunks\n"
        "        │\n"
        "        ▼\n"
        "LLMService.generate()        → httpx async POST\n"
        "        │                       Ollama :11434  (model = llama3)\n"
        "        ▼\n"
        "ChatResponse{ answer, sources[], brand_logo, language=\"ko\" }"
    )
    code_block(s, Inches(0.6), Inches(2.0), Inches(8.0), Inches(4.0), flow)
    # right column
    card(s, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.0))
    tx(s, Inches(9.15), Inches(2.15), Inches(3.4), Inches(0.4),
       "WHY LOCAL", size=11, bold=True, color=TEAL, font=FONT)
    notes = [
        "• Zero per-call cost — Ollama on the box.",
        "• PII / health context never leaves the host.",
        "• Swap llama3 → HyperCLOVA-X without API change.",
        "• SSE streaming via /chat/stream for progressive UX.",
    ]
    yy = Inches(2.6)
    for n in notes:
        tx(s, Inches(9.15), yy, Inches(3.4), Inches(0.7),
           n, size=11, color=INK, font=FONT)
        yy += Inches(0.7)
    footer(s, page, total)


def s_persona(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · Persona engine", "해요체 + 받침/조사 auto-resolution")
    code = (
        "# templates.py — auto-pick 은/는, 이/가, 을/를 by 받침\n"
        "def josa(word: str, has_batchim: str, no_batchim: str) -> str:\n"
        "    last = word[-1]\n"
        "    code = (ord(last) - 0xAC00) % 28 if 0xAC00 <= ord(last) <= 0xD7A3 else 0\n"
        "    return has_batchim if code else no_batchim\n"
        "\n"
        "# usage\n"
        "msg = f\"{studio}{josa(studio, '은', '는')} 지금 수련 가능해요.\"\n"
        "# → \"성수 elbee 요가 스튜디오는 지금 수련 가능해요.\""
    )
    code_block(s, Inches(0.6), Inches(2.0), Inches(8.0), Inches(2.6), code)
    card(s, Inches(8.9), Inches(2.0), Inches(3.8), Inches(2.6))
    tx(s, Inches(9.15), Inches(2.15), Inches(3.4), Inches(0.4),
       "ELBEE PERSONA", size=11, bold=True, color=GOLD, font=FONT)
    tx(s, Inches(9.15), Inches(2.55), Inches(3.4), Inches(2.0),
       "• 해요체 default register\n"
       "• Brand: elbee.yogaman.club\n"
       "• Proximity messages\n"
       "• Weather-aware routing\n"
       "• Few-shot tone-and-manner",
       size=11, color=INK, font=FONT)
    # bottom: sample output
    card(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(1.85))
    tx(s, Inches(0.85), Inches(4.95), Inches(11.6), Inches(0.4),
       "SAMPLE invite_message (LocationSearchResponse)",
       size=11, bold=True, color=TEAL, font=FONT)
    tx(s, Inches(0.85), Inches(5.4), Inches(11.6), Inches(1.3),
       "안녕하세요, elbee 멤버! 🙏 오늘 성수동은 근처 성수 elbee 요가 스튜디오에서\n"
       "수련하실 기회가 생겼어요. 우천 시 실내 수업으로 자동 매칭됩니다.",
       size=14, color=INK, font=FONT)
    footer(s, page, total)


def s_api_surface(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · API surface", "11 routes · OpenAPI · Pydantic v2")
    rows = [
        ("POST", "/chat",                "RAG chat — 해요체 answer + sources"),
        ("POST", "/chat/stream",         "SSE streaming variant"),
        ("GET",  "/chat/health",         "LLM + data store liveness"),
        ("POST", "/search",              "Keyword search over 165 pages"),
        ("GET",  "/search/suggest",      "Prefix autocomplete"),
        ("GET",  "/search/filters",      "Available platforms · regions · topics"),
        ("GET",  "/search/health",       "Index size + readiness"),
        ("POST", "/search/locations",    "Yoga spot proximity / amenity / weather"),
        ("GET",  "/docs",                "Swagger UI"),
        ("GET",  "/redoc",               "ReDoc UI"),
        ("GET",  "/health",              "App-level health"),
    ]
    # header row
    yy = Inches(2.0)
    tx(s, Inches(0.6),  yy, Inches(0.9), Inches(0.3), "METHOD", size=10, bold=True, color=MUTED, font=FONT)
    tx(s, Inches(1.6),  yy, Inches(3.5), Inches(0.3), "PATH",   size=10, bold=True, color=MUTED, font=FONT)
    tx(s, Inches(5.2),  yy, Inches(7.5), Inches(0.3), "DESCRIPTION", size=10, bold=True, color=MUTED, font=FONT)
    yy += Inches(0.35)
    for m, p, d in rows:
        bar(s, Inches(0.6), yy + Inches(0.05), w=Inches(12.1), h=Inches(0.01), color=DIM)
        accent = TEAL if m == "GET" else GOLD
        tx(s, Inches(0.6), yy + Inches(0.08), Inches(0.9), Inches(0.32),
           m, size=11, bold=True, color=accent, font=FONT)
        tx(s, Inches(1.6), yy + Inches(0.08), Inches(3.5), Inches(0.32),
           p, size=11, color=INK, font="Consolas")
        tx(s, Inches(5.2), yy + Inches(0.08), Inches(7.5), Inches(0.32),
           d, size=11, color=MUTED, font=FONT)
        yy += Inches(0.4)
    footer(s, page, total)


def s_models(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · Schemas", "Pydantic v2 — typed, validated, OpenAPI-published")
    chat = (
        "ChatResponse {\n"
        "  answer: str,\n"
        "  sources: [{ page: int, snippet: str, topic: str }],\n"
        "  model: 'llama3',\n"
        "  brand_logo: HttpUrl,\n"
        "  language: 'ko'\n"
        "}"
    )
    loc = (
        "LocationSearchResponse {\n"
        "  results: [YogaLocationResult{\n"
        "    id, name, type, district, address, lat, lng,\n"
        "    amenities[], tags[], weather_indoor: bool,\n"
        "    rating: float, distance_km: float,\n"
        "    invite_message: str   # 해요체\n"
        "  }],\n"
        "  total, message, brand, brand_logo\n"
        "}"
    )
    code_block(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.7), chat)
    code_block(s, Inches(6.7), Inches(2.0), Inches(6.0), Inches(4.7), loc)
    footer(s, page, total)


def s_twin(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · Cross-platform twin", "macOS dev → WSL2 / Win11 parity")
    rows = [
        ("Item",            "macOS",                              "Windows / WSL2"),
        ("Tesseract",       "/opt/homebrew/bin/tesseract",        "C:\\Program Files\\Tesseract-OCR\\tesseract.exe"),
        ("Ollama endpoint", "http://localhost:11434",             "http://host.docker.internal:11434"),
        ("Line endings",    "LF",                                  "git config core.autocrlf input"),
        ("Python",          "python3",                             "python  (or python3 in WSL2)"),
        ("Bootstrap",       "brew install tesseract ollama",       "wsl --install -d Ubuntu-22.04"),
    ]
    yy = Inches(2.0)
    for i, (a, b, c) in enumerate(rows):
        is_head = (i == 0)
        col_a = MUTED if is_head else TEAL
        col_b = MUTED if is_head else INK
        col_c = MUTED if is_head else INK
        bar(s, Inches(0.6), yy + Inches(0.05), w=Inches(12.1), h=Inches(0.01), color=DIM)
        tx(s, Inches(0.6),  yy + Inches(0.08), Inches(2.4), Inches(0.32),
           a, size=11, bold=True, color=col_a, font=FONT)
        tx(s, Inches(3.0),  yy + Inches(0.08), Inches(4.7), Inches(0.32),
           b, size=11, color=col_b, font="Consolas" if not is_head else FONT)
        tx(s, Inches(7.8),  yy + Inches(0.08), Inches(4.9), Inches(0.32),
           c, size=11, color=col_c, font="Consolas" if not is_head else FONT)
        yy += Inches(0.55)
    footer(s, page, total)


def s_roadmap(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    header(s, f"{page:02d} · Roadmap & open issues", "What's shipped, what's next")
    phases = [
        ("Phase 0", "Scaffold + OCR pipeline",                     "DONE",  TEAL),
        ("Phase 1", "165-page intake (≈28k words indexed)",         "DONE",  TEAL),
        ("Phase 2", "elbee.yogaman.club brand · 8 yoga spots",      "DONE",  TEAL),
        ("Phase 3", "Multi-book expansion + manifest",              "NEXT",  GOLD),
        ("Phase 4", "GitHub Actions CI (OCR lint + JSON schema)",   "PLAN",  MUTED),
        ("Phase 5", "Employer demo package + live deploy",          "PLAN",  MUTED),
    ]
    yy = Inches(2.0)
    for ph, body, tag, col in phases:
        card(s, Inches(0.6), yy, Inches(8.0), Inches(0.65))
        tx(s, Inches(0.85), yy + Inches(0.16), Inches(1.4), Inches(0.4),
           ph, size=13, bold=True, color=col, font=FONT)
        tx(s, Inches(2.3), yy + Inches(0.16), Inches(5.5), Inches(0.4),
           body, size=12, color=INK, font=FONT)
        chip(s, Inches(7.6), yy + Inches(0.18), tag, color=col)
        yy += Inches(0.78)
    # right: open issues
    card(s, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.7))
    tx(s, Inches(9.15), Inches(2.15), Inches(3.4), Inches(0.4),
       "OPEN ISSUES", size=11, bold=True, color=GOLD, font=FONT)
    issues = [
        "#1  SSOT architecture (this deck's source)",
        "#3  60s vertical storyboard",
        "#4  Korean NLP + GEO techstack",
    ]
    yy = Inches(2.6)
    for n in issues:
        tx(s, Inches(9.15), yy, Inches(3.4), Inches(0.6),
           "• " + n, size=11, color=INK, font=FONT)
        yy += Inches(0.6)
    footer(s, page, total)


def s_thanks(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               Inches(0.18), H)
    block.line.fill.background(); block.fill.solid()
    block.fill.fore_color.rgb = TEAL
    tx(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.4),
       "Built local-first.", size=58, bold=True, color=INK, font=FONT)
    tx(s, Inches(0.8), Inches(3.7), Inches(11.5), Inches(1.0),
       "Korean-native. Cited by AI. Ready to scale.",
       size=22, color=MUTED, font=FONT)
    chip(s, Inches(0.8), Inches(5.4), "github.com/aiegoo/aeogeo")
    chip(s, Inches(4.0), Inches(5.4), "elbee.yogaman.club", color=GOLD)
    footer(s, page, total)


def main():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    builders = [
        s_title,
        s_system_map,
        s_korean_nlp,
        s_geo_context,
        s_ocr_pipeline,
        s_stores,
        s_rag_llm,
        s_persona,
        s_api_surface,
        s_models,
        s_twin,
        s_roadmap,
        s_thanks,
    ]
    total = len(builders)
    s_title(prs)
    for i, fn in enumerate(builders[1:-1], start=2):
        fn(prs, i, total)
    s_thanks(prs, total, total)
    out = Path(__file__).parent / "aeogeo_tech_deepdive_deck.pptx"
    prs.save(out)
    print(f"Wrote {out} {out.stat().st_size} bytes, {total} slides")


if __name__ == "__main__":
    main()
