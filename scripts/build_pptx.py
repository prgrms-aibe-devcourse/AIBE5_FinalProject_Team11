#!/usr/bin/env python3
"""
build_pptx.py — Generate 요가큐 presentation as a PowerPoint file.
Saves to /mnt/c/Users/hsyyu/Downloads/요가큐_발표.pptx (Windows Downloads via WSL2).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x1a, 0x1a, 0x2e)
BG_MID       = RGBColor(0x1e, 0x1b, 0x4b)
ACCENT_VIOL  = RGBColor(0xa7, 0x8b, 0xfa)
ACCENT_SKY   = RGBColor(0x7d, 0xd3, 0xfc)
ACCENT_GREEN = RGBColor(0x86, 0xef, 0xac)
ACCENT_AMBER = RGBColor(0xfb, 0xbf, 0x24)
TEXT_MAIN    = RGBColor(0xe2, 0xe8, 0xf0)
TEXT_DIM     = RGBColor(0x94, 0xa3, 0xb8)
TH_BG        = RGBColor(0x31, 0x2e, 0x81)
TD_BG        = RGBColor(0x1e, 0x1b, 0x4b)
TD_ALT       = RGBColor(0x25, 0x22, 0x5e)
RED_WARN     = RGBColor(0xf8, 0x71, 0x71)

# ── Slide size: 16:9 widescreen ───────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # completely blank layout


# ── Helpers ───────────────────────────────────────────────────────────────────
def bg(slide, color=BG_DARK):
    """Fill slide background with solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txbox(slide, text, left, top, width, height,
          size=20, bold=False, color=TEXT_MAIN, align=PP_ALIGN.LEFT,
          italic=False, wrap=True):
    """Add a text box and return the shape."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_bar(slide, title, subtitle=None):
    """Render a standard h2-style title bar at top of slide."""
    txbox(slide, title,
          left=Inches(0.5), top=Inches(0.3),
          width=Inches(12.3), height=Inches(0.7),
          size=28, bold=True, color=ACCENT_SKY)
    # underline rule
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    # draw a thin rectangle as underline
    rule = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), Inches(1.0), Inches(12.3), Pt(2)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor(0x6C, 0x63, 0xFF)
    rule.line.fill.background()

    if subtitle:
        txbox(slide, subtitle,
              left=Inches(0.5), top=Inches(1.1),
              width=Inches(12.3), height=Inches(0.4),
              size=16, color=TEXT_DIM)


def bullet_block(slide, lines, top, size=19, indent=0):
    """Render a list of (text, color, bold) tuples as stacked text boxes."""
    y = top
    left  = Inches(0.5 + indent)
    width = Inches(12.3 - indent)
    for (text, color, bold) in lines:
        h = Inches(0.45)
        txbox(slide, text, left=left, top=y, width=width, height=h,
              size=size, bold=bold, color=color)
        y += h
    return y


def table_shape(slide, headers, rows, top, col_widths=None):
    """Draw a styled table."""
    ncols = len(headers)
    nrows = len(rows)
    if col_widths is None:
        col_widths = [Inches(12.3 / ncols)] * ncols

    left = Inches(0.5)
    row_h = Inches(0.42)
    tbl = slide.shapes.add_table(
        nrows + 1, ncols,
        left, top,
        sum(col_widths), row_h * (nrows + 1)
    ).table

    # set column widths
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    # header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = TH_BG
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.text = hdr
        run.font.bold  = True
        run.font.size  = Pt(16)
        run.font.color.rgb = ACCENT_VIOL

    # data rows
    for ri, row in enumerate(rows):
        bg_c = TD_ALT if ri % 2 else TD_BG
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_c
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.text = val
            run.font.size  = Pt(15)
            run.font.color.rgb = TEXT_MAIN

    return tbl


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Intro (title slide)
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)

txbox(sl, "요가큐",
      Inches(1), Inches(1.5), Inches(11.3), Inches(1.6),
      size=72, bold=True, color=ACCENT_VIOL, align=PP_ALIGN.CENTER)

txbox(sl, "AI 기반 요가 포즈 매칭 & 안전 필터 플랫폼",
      Inches(1), Inches(3.2), Inches(11.3), Inches(0.7),
      size=28, color=TEXT_MAIN, align=PP_ALIGN.CENTER)

txbox(sl, "match.yogaman.club  |  elbee.yogaman.club",
      Inches(1), Inches(4.0), Inches(11.3), Inches(0.5),
      size=20, color=ACCENT_SKY, align=PP_ALIGN.CENTER)

txbox(sl, "수련자의 신체 조건과 목표에 맞는 안전하고 정확한 포즈 추천 시스템",
      Inches(1), Inches(4.7), Inches(11.3), Inches(0.6),
      size=20, color=TEXT_DIM, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — 발표 순서
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
title_bar(sl, "발표 순서")

items = [
    "1.  프로젝트 배경 — 문제 정의",
    "2.  기존 플랫폼 분석 — 검색·소셜 미디어 / 예약 플랫폼",
    "3.  차별점 — 5대 한계와 해결 방향",
    "4.  활용 데이터 — 포즈 DB · Q&A · OCR",
    "5.  기술 스택 — Spring Boot · FastAPI · RAG",
    "6.  주요 기능 시연",
    "7.  트러블슈팅 — 4가지 핵심 문제 해결",
    "8.  향후 개선 방향",
]
y = Inches(1.3)
for item in items:
    txbox(sl, item, Inches(1.2), y, Inches(11), Inches(0.48),
          size=21, color=TEXT_MAIN)
    y += Inches(0.52)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — 프로젝트 배경
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
title_bar(sl, "프로젝트 배경")

txbox(sl, "수련자의 안전 공백",
      Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.45),
      size=22, bold=True, color=ACCENT_GREEN)

body1 = [
    "• 허리 디스크·무릎 부상·고혈압 — 검증되지 않은 콘텐츠로 부상 악화",
    '• "초보자 요가" 검색 → 척추 부담 포즈가 상위 노출',
]
y = Inches(1.7)
for b in body1:
    txbox(sl, b, Inches(0.8), y, Inches(11.5), Inches(0.45), size=19, color=TEXT_MAIN)
    y += Inches(0.48)

txbox(sl, "강사의 디지털 가시성 부재",
      Inches(0.5), Inches(2.75), Inches(12.3), Inches(0.45),
      size=22, bold=True, color=ACCENT_GREEN)

body2 = [
    "• 인스타그램 강사 프로필 포화 — 단편 게시물만으로 전문성·자격 구별 불가",
    "• 20년 경력 강사도 파편화된 SNS 속에 묻혀 신뢰 신호 전달 못함",
    "• AI 시대 검색·추천 흐름에서 지속적으로 이탈",
]
y = Inches(3.2)
for b in body2:
    txbox(sl, b, Inches(0.8), y, Inches(11.5), Inches(0.45), size=19, color=TEXT_MAIN)
    y += Inches(0.48)

txbox(sl, "➜ 수련자 안전 공백 + 강사 디지털 가시성 부재를 동시에 해결하는 신뢰 기반 AI 매칭 플랫폼",
      Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.6),
      size=19, bold=True, color=ACCENT_AMBER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — 기존 플랫폼 분석 ①
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
title_bar(sl, "기존 플랫폼 분석 ①", "검색·소셜 미디어 채널 (네이버·구글·인스타그램·유튜브)")

table_shape(
    sl,
    headers=["채널", "장점", "한계"],
    rows=[
        ["네이버·구글",  "콘텐츠 양, 접근성",            "상위 = 광고·블로그, 신체 조건 미반영"],
        ["유튜브",       "풍부한 영상 콘텐츠",            "조회수 중심 알고리즘 → 자격 무관 노출"],
        ["인스타그램",   "실시간 강사 활동 확인",         "프로필 포화 → 전문성 구별 불가"],
    ],
    top=Inches(1.5),
    col_widths=[Inches(2.8), Inches(4.2), Inches(5.3)],
)

txbox(sl, "공통 문제: 금기 포즈 필터링 기능 전혀 없음 → 부상자에게 오히려 위험한 콘텐츠 추천",
      Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.55),
      size=20, bold=True, color=RED_WARN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — 기존 플랫폼 분석 ②
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
title_bar(sl, "기존 플랫폼 분석 ②", "예약 마켓플레이스 — Mindbody · ClassPass")

txbox(sl, "✅  장점",
      Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.45),
      size=22, bold=True, color=ACCENT_GREEN)
pros = [
    "• 스튜디오 검색·예약 편의성 높음",
    "• 글로벌 네트워크 보유",
]
y = Inches(1.95)
for p in pros:
    txbox(sl, p, Inches(0.9), y, Inches(11.5), Inches(0.42), size=19, color=TEXT_MAIN)
    y += Inches(0.45)

txbox(sl, "❌  한계",
      Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.45),
      size=22, bold=True, color=RED_WARN)
cons = [
    "• 단순 예약 인터페이스 — 건강 상태·목표 기반 추천 없음",
    "• Schema.org 구조화 데이터 미제공",
    "  → Gemini · SearchGPT · Perplexity 등 AI 검색엔진 인용 불가",
]
y = Inches(3.45)
for c in cons:
    txbox(sl, c, Inches(0.9), y, Inches(11.5), Inches(0.42), size=19, color=TEXT_MAIN)
    y += Inches(0.45)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — 차별점
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
title_bar(sl, "차별점 — 5대 한계 해결")

table_shape(
    sl,
    headers=["기존 문제", "요가큐 해결 방식"],
    rows=[
        ["안전성 판단 불가",  "Kill-Switch 금기사항 하드 차단 — 위험 포즈 자동 제외"],
        ["강사 자격 불투명",  "instructor_trust_score + Schema.org Person JSON-LD"],
        ["개인화 필터 없음",  "경험 레벨 · 가용 시간 · 목표 다중 필터 매칭"],
        ["AI 인용 불가",      "FAQPage · HowTo · DefinedTerm JSON-LD 자동 생성"],
        ["위치 매칭 없음",    "Haversine 반경 쿼리 + 전문 분야 가중치 랭킹"],
    ],
    top=Inches(1.3),
    col_widths=[Inches(3.5), Inches(8.8)],
)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — 활용 데이터
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
title_bar(sl, "활용 데이터")

table_shape(
    sl,
    headers=["데이터", "규모", "활용"],
    rows=[
        ["요가 포즈 DB",      "2,700+ 아사나",   "E-E-A-T 메타데이터 enrichment → PostgreSQL"],
        ["Q&A 학습 데이터",   "2,899 FAQ 쌍",    "GPT-4o-mini 생성 → FAQPage JSON-LD"],
        ["강사·스튜디오",     "서울 10개소",     "위치 기반 매칭 + Person JSON-LD"],
        ["OCR 교육 자료",     "전문 도서",       "Tesseract → 벡터 임베딩 → RAG 챗봇 인용"],
    ],
    top=Inches(1.4),
    col_widths=[Inches(3.0), Inches(2.8), Inches(6.5)],
)

txbox(sl, "856개 포즈 자연어 설명  ·  nomic-embed-text 임베딩  ·  mistral:latest 로컬 LLM",
      Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.5),
      size=17, color=TEXT_DIM, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — 기술 스택
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
title_bar(sl, "기술 스택")

stack = [
    ("Backend",    "Spring Boot 3.x  ·  Java 17  ·  Spring Data JPA  ·  Flyway  ·  PostgreSQL 16"),
    ("AI Layer",   "FastAPI  ·  Ollama (mistral · nomic-embed-text)  ·  GPT-4o-mini  ·  LangChain RAG"),
    ("Frontend",   "React 18  ·  Vite 5  ·  Streamlit (매칭 점수 시각화 · 안전 필터 데모)"),
    ("Infra",      "Docker Compose  ·  Caddy  ·  Cloudflare Tunnel  ·  GitHub Actions CI"),
    ("API Docs",   "SpringDoc OpenAPI  →  match.yogaman.club/swagger-ui.html"),
]
y = Inches(1.35)
for label, content in stack:
    txbox(sl, label,
          Inches(0.5), y, Inches(2.5), Inches(0.5),
          size=18, bold=True, color=ACCENT_SKY)
    txbox(sl, content,
          Inches(3.0), y, Inches(9.8), Inches(0.5),
          size=17, color=TEXT_MAIN)
    y += Inches(0.8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — 주요 기능
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
title_bar(sl, "주요 기능")

txbox(sl, "수련자",
      Inches(0.5), Inches(1.3), Inches(6.0), Inches(0.45),
      size=21, bold=True, color=ACCENT_GREEN)
user_items = [
    "• 자연어 입력 → AI 포즈 매칭 (/api/v1/match)",
    "• Kill-Switch: 금기 포즈 자동 차단 + 사유 반환",
    "• 경험 레벨 · 가용 시간 필터",
    "• 위치 기반 강사·스튜디오 추천",
    "• RAG 챗봇 자유 질문 (/chat)",
]
y = Inches(1.75)
for item in user_items:
    txbox(sl, item, Inches(0.8), y, Inches(5.8), Inches(0.43), size=17, color=TEXT_MAIN)
    y += Inches(0.46)

txbox(sl, "강사·스튜디오",
      Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.45),
      size=21, bold=True, color=ACCENT_AMBER)
biz_items = [
    "• Schema.org Person JSON-LD 자동 생성",
    "  (/instructors/{id}/jsonld)",
    "• 포즈별 FAQPage · HowTo · DefinedTerm JSON-LD",
    "• Google 리치 결과 및 Perplexity 인용 확보",
    "• 아티클 콘텐츠 관리 (/articles)",
]
y = Inches(1.75)
for item in biz_items:
    txbox(sl, item, Inches(7.0), y, Inches(5.8), Inches(0.43), size=17, color=TEXT_MAIN)
    y += Inches(0.46)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — 시연
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)

txbox(sl, "🎬  시연",
      Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
      size=60, bold=True, color=ACCENT_VIOL, align=PP_ALIGN.CENTER)
txbox(sl, "match.yogaman.club",
      Inches(1), Inches(3.6), Inches(11.3), Inches(0.6),
      size=26, color=ACCENT_SKY, align=PP_ALIGN.CENTER)
txbox(sl, "매칭 시연  →  RAG 챗봇  →  JSON-LD AEO  →  안전 필터",
      Inches(1), Inches(4.3), Inches(11.3), Inches(0.5),
      size=20, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — 트러블슈팅
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
title_bar(sl, "트러블슈팅")

troubles = [
    ("① 매칭 점수 항상 0.0",
     "원인: UI 목표값(Spinal_Mobility) ≠ DB 태그 어휘(mobility, back)",
     "해결: GOAL_TAG_MAP — UI 목표 → DB 태그 확장 매핑 도입"),
    ("② Kill-Switch 하드 차단 미작동",
     "원인: 단순 조건 확인으로는 severity 단계 차단 불가",
     "해결: score=0, blocked=true 고정 반환 + 차단 사유 명시"),
    ("③ JSON-LD @type 필드 누락",
     "원인: Jackson null 필드 직렬화 예외",
     "해결: NON_NULL 설정 + LinkedHashMap 명시적 구성"),
    ("④ Streamlit 컨테이너 404",
     "원인: --server.baseUrlPath /safety → 루트 / 미응답",
     "해결: baseUrlPath 제거 + Caddyfile /safety* → 8501 라우팅"),
]

y = Inches(1.35)
for (title, cause, fix) in troubles:
    txbox(sl, title, Inches(0.5), y, Inches(12.3), Inches(0.38),
          size=18, bold=True, color=ACCENT_AMBER)
    y += Inches(0.38)
    txbox(sl, cause, Inches(0.9), y, Inches(12.0), Inches(0.35),
          size=16, color=TEXT_DIM)
    y += Inches(0.35)
    txbox(sl, fix, Inches(0.9), y, Inches(12.0), Inches(0.35),
          size=16, color=ACCENT_GREEN)
    y += Inches(0.48)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — 향후 개선 방향 + 감사합니다
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
title_bar(sl, "향후 개선 방향")

roadmap = [
    ("1. 예약 시스템 통합",
     "Mindbody·1club API → 매칭 → 예약 전환 풀 파이프라인"),
    ("2. AI 인용 범위 확대",
     "전체 2,700 포즈 JSON-LD → Perplexity·Gemini 인용 빈도 확대"),
    ("3. 수련 이력 기반 개인화",
     "세션 이력 축적 → 점진적 난이도 조절 · 목표 추적"),
    ("4. 모바일 앱 전환",
     "React Native 확장 → 더 넓은 사용자층 도달"),
]
y = Inches(1.3)
for (head, body) in roadmap:
    txbox(sl, head, Inches(0.5), y, Inches(12.3), Inches(0.42),
          size=20, bold=True, color=ACCENT_SKY)
    y += Inches(0.42)
    txbox(sl, body, Inches(0.9), y, Inches(12.0), Inches(0.38),
          size=17, color=TEXT_MAIN)
    y += Inches(0.52)

txbox(sl, "감사합니다",
      Inches(1), Inches(6.2), Inches(11.3), Inches(0.8),
      size=36, bold=True, color=ACCENT_VIOL, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "/mnt/c/Users/hsyyu/Downloads/요가큐_발표.pptx"
prs.save(OUT)
print(f"✅  Saved → {OUT}")
