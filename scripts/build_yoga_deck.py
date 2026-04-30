#!/usr/bin/env python3
"""
build_yoga_deck.py
──────────────────
Generates 요가큐 investor/judge presentation with:
  - Yoga/restorative warm-earth dark theme
  - Full condition-check → metadata mapping → provider data → user journey flow
  - 4 user capabilities: instructor / shorts / studio / HUD cueing

Output: /mnt/c/Users/hsyyu/Downloads/요가큐_신규발표.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

# ── Yoga / Restorative Earth Palette ─────────────────────────────────────────
BG_DEEP    = RGBColor(0x0F, 0x17, 0x0B)   # deep forest night
BG_CARD    = RGBColor(0x1A, 0x24, 0x14)   # dark sage card
BG_CARD2   = RGBColor(0x1F, 0x2B, 0x18)   # slightly lighter card
SAGE       = RGBColor(0x7C, 0xB8, 0x7A)   # sage green (main accent)
EARTH      = RGBColor(0xC4, 0x86, 0x6A)   # terracotta / earth
GOLD       = RGBColor(0xD4, 0xA8, 0x53)   # warm gold
LOTUS      = RGBColor(0xB5, 0xA4, 0xC9)   # soft lotus lavender
TEXT_MAIN  = RGBColor(0xF0, 0xED, 0xE6)   # warm white
TEXT_DIM   = RGBColor(0x9F, 0xAF, 0x8C)   # muted sage
TEXT_GOLD  = RGBColor(0xD4, 0xA8, 0x53)
KILL_RED   = RGBColor(0xE0, 0x70, 0x70)   # soft kill-switch red
RULE_COLOR = RGBColor(0x4A, 0x7C, 0x3F)   # dark sage rule line
HEADER_BG  = RGBColor(0x25, 0x36, 0x1C)   # table header bg
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


# ── Core helpers ──────────────────────────────────────────────────────────────

def bg(slide, color=BG_DEEP):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill_color, line_color=None, line_width_pt=0):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = _Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, text, left, top, width, height,
          size=18, bold=False, italic=False, color=TEXT_MAIN,
          align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def multiline_txbox(slide, lines, left, top, width, height, default_size=18, wrap=True):
    """lines: list of (text, size, bold, color, align)"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    first = True
    for (text, size, bold, color, align) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def title_bar(slide, title, subtitle=None, title_color=SAGE):
    txbox(slide, title,
          Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.7),
          size=30, bold=True, color=title_color)
    # sage rule
    rect(slide, Inches(0.55), Inches(0.95), Inches(12.2), Pt(2.5), RULE_COLOR)
    if subtitle:
        txbox(slide, subtitle,
              Inches(0.55), Inches(1.0), Inches(12.2), Inches(0.45),
              size=15, color=TEXT_DIM, italic=True)


def card(slide, left, top, width, height, fill=BG_CARD, border=RULE_COLOR, border_pt=1):
    return rect(slide, left, top, width, height, fill, border, border_pt)


def badge(slide, text, left, top, width=Inches(2.2), height=Inches(0.38),
          bg_color=RULE_COLOR, text_color=TEXT_MAIN, size=13):
    rect(slide, left, top, width, height, bg_color, line_color=None)
    txbox(slide, text, left, top, width, height,
          size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def icon_card(slide, icon, label, sublabel, left, top, width=Inches(2.9), height=Inches(1.8),
              accent=SAGE):
    card(slide, left, top, width, height, BG_CARD2, accent, 1)
    txbox(slide, icon, left, top + Inches(0.18), width, Inches(0.55),
          size=28, align=PP_ALIGN.CENTER, color=accent)
    txbox(slide, label, left, top + Inches(0.72), width, Inches(0.45),
          size=15, bold=True, align=PP_ALIGN.CENTER, color=TEXT_MAIN)
    txbox(slide, sublabel, left, top + Inches(1.15), width, Inches(0.55),
          size=11, align=PP_ALIGN.CENTER, color=TEXT_DIM, wrap=True)


def flow_arrow(slide, left, top):
    """A small right-pointing arrow between flow nodes."""
    rect(slide, left, top + Inches(0.12), Inches(0.22), Pt(3), RULE_COLOR)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════════════════════════════════

def slide_cover():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    # left accent bar
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, SAGE)
    # large title
    txbox(sl, "요가큐 · Yogaman.club",
          Inches(0.5), Inches(1.8), Inches(8), Inches(1.1),
          size=44, bold=True, color=TEXT_MAIN)
    txbox(sl, "AI 기반 요가 매칭 플랫폼",
          Inches(0.5), Inches(2.9), Inches(8), Inches(0.65),
          size=26, bold=False, color=SAGE)
    txbox(sl,
          "당신의 몸 상태를 이해하고\n안전한 포즈 · 강사 · 스튜디오 · 홈 큐잉을 연결합니다.",
          Inches(0.5), Inches(3.65), Inches(8.5), Inches(1.0),
          size=18, color=TEXT_DIM, wrap=True)
    # right decorative block
    rect(sl, Inches(9.5), Inches(1.4), Inches(3.3), Inches(4.6), BG_CARD, RULE_COLOR, 1)
    lines = [
        ("🧘  몸 상태 체크", 15, False, SAGE),
        ("→  Kill-Switch 안전 필터", 14, False, KILL_RED),
        ("→  포즈 / 스튜디오 매칭", 14, False, TEXT_MAIN),
        ("→  강사 · 쇼츠 발견", 14, False, EARTH),
        ("→  홈 큐잉 + 보이스오버", 14, False, LOTUS),
    ]
    y = Inches(1.75)
    for (t, s, b, c) in lines:
        txbox(sl, t, Inches(9.7), y, Inches(2.9), Inches(0.45),
              size=s, bold=b, color=c)
        y += Inches(0.58)
    # domain badges
    for i, domain in enumerate(["match.yogaman.club", "elbee.yogaman.club"]):
        badge(sl, domain, Inches(0.5 + i * 3.0), Inches(6.8),
              width=Inches(2.6), bg_color=HEADER_BG, text_color=GOLD, size=12)
    txbox(sl, "AIBE5 · Team 14", Inches(11.0), Inches(7.1),
          Inches(2.2), Inches(0.35), size=12, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem: 기존 앱의 한계
# ═══════════════════════════════════════════════════════════════════════════════

def slide_problem():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, KILL_RED)
    title_bar(sl, "01 · 문제 — 기존 앱은 당신의 몸을 모릅니다", title_color=KILL_RED)

    # 3 pain-point cards
    pains = [
        ("⚠️  금기 포즈 필터 없음",
         "허리 디스크 환자에게도 백벤드 추천.\n유튜브·네이버 알고리즘은 조회수만 최적화.",
         KILL_RED),
        ("📋  강사 검증 불가",
         "자격증 유무 표시 없음.\nMindbody·ClassPass도 E-E-A-T 신호 부재.",
         EARTH),
        ("🗺️  스튜디오 = 리스트뷰",
         "거리·전문화·인증 복합 매칭 없음.\n예약까지 3~5 앱을 전전.",
         LOTUS),
    ]
    for i, (title, body, accent) in enumerate(pains):
        x = Inches(0.55 + i * 4.2)
        card(sl, x, Inches(1.55), Inches(3.9), Inches(2.4), BG_CARD, accent, 1)
        txbox(sl, title, x + Inches(0.15), Inches(1.7), Inches(3.6), Inches(0.5),
              size=16, bold=True, color=accent)
        txbox(sl, body, x + Inches(0.15), Inches(2.2), Inches(3.6), Inches(1.5),
              size=13, color=TEXT_DIM, wrap=True)

    # quote
    rect(sl, Inches(0.55), Inches(4.2), Inches(12.2), Inches(1.1), BG_CARD2, RULE_COLOR, 1)
    txbox(sl,
          '"잘못된 자세 하나가 디스크 수술로 이어질 수 있습니다. 그래서 우리는 안전을 먼저 설계했습니다."',
          Inches(0.8), Inches(4.3), Inches(11.8), Inches(0.9),
          size=17, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

    # gap stats row
    stats = [("1,150억 $", "글로벌 요가 시장 (2025)"), ("68%", "앱 이탈 — 맞춤 추천 부재"), ("0", "부상자용 안전 필터 앱")]
    for i, (num, label) in enumerate(stats):
        x = Inches(1.2 + i * 3.9)
        txbox(sl, num, x, Inches(5.5), Inches(3.5), Inches(0.7),
              size=32, bold=True, color=SAGE, align=PP_ALIGN.CENTER)
        txbox(sl, label, x, Inches(6.15), Inches(3.5), Inches(0.45),
              size=13, color=TEXT_DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 몸 상태 체크 → 메타데이터 매핑 (Condition Check)
# ═══════════════════════════════════════════════════════════════════════════════

def slide_condition_check():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, SAGE)
    title_bar(sl, "02 · 몸 상태 체크 → 메타데이터 매핑")

    # Step flow: User Input → Condition Tags → Kill-Switch → Safe Candidates
    flow_items = [
        ("👤  사용자 입력",   "• 허리 디스크\n• 고혈압\n• 임신 중\n• 무릎 수술 후", SAGE),
        ("🔖  컨디션 태그",   "• disc / hypertension\n• pregnancy\n• knee_surgery\n태그 → DB 조회", EARTH),
        ("⛔  Kill-Switch",   "kill_switch = TRUE\nCRITICAL 등급\n→ 완전 제외\n점수 무관", KILL_RED),
        ("✅  안전 후보군",   "CAUTION 등급:\n수정 큐잉 제공\nMEDICAL_CLEARANCE:\n의사 확인 권고", LOTUS),
        ("🏆  매칭 & 순위",   "benefit.weight\n× goal_overlap\n→ Top-K 반환", GOLD),
    ]
    box_w = Inches(2.35)
    arrow_w = Inches(0.22)
    start_x = Inches(0.4)
    y = Inches(1.6)
    h = Inches(3.4)

    for i, (title, body, accent) in enumerate(flow_items):
        x = start_x + i * (box_w + arrow_w)
        card(sl, x, y, box_w, h, BG_CARD, accent, 1)
        txbox(sl, title, x + Inches(0.1), y + Inches(0.15), box_w - Inches(0.2), Inches(0.5),
              size=14, bold=True, color=accent)
        txbox(sl, body, x + Inches(0.1), y + Inches(0.65), box_w - Inches(0.2), h - Inches(0.8),
              size=12, color=TEXT_MAIN, wrap=True)
        # arrow between boxes
        if i < len(flow_items) - 1:
            ax = x + box_w + Inches(0.02)
            ay = y + h / 2 - Inches(0.1)
            rect(sl, ax, ay, Inches(0.18), Pt(3), TEXT_DIM)

    # Schema strip at bottom
    rect(sl, Inches(0.4), Inches(5.2), Inches(12.5), Inches(1.05), BG_CARD2, RULE_COLOR, 1)
    txbox(sl, "📐  메타데이터 스키마 (pose_eat_schema.json)",
          Inches(0.6), Inches(5.28), Inches(12.0), Inches(0.38),
          size=13, bold=True, color=SAGE)
    schema_text = (
        'contraindications[]:  { condition, severity: "CAUTION | CRITICAL | MEDICAL_CLEARANCE", '
        'kill_switch: bool, instruction }    ·    '
        'benefits[]:  { tag, weight: 0–1 }    ·    '
        'difficulty_rank: 1–5    ·    Schema.org ExerciseAction JSON-LD'
    )
    txbox(sl, schema_text, Inches(0.6), Inches(5.65), Inches(12.0), Inches(0.52),
          size=11, color=TEXT_DIM, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 프로바이더 데이터 수집 (Provider Data)
# ═══════════════════════════════════════════════════════════════════════════════

def slide_provider_data():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, EARTH)
    title_bar(sl, "03 · 프로바이더 데이터 — 매칭 엔진을 작동시키는 원천", title_color=EARTH)

    # 3-column: Studio | Instructor | Content
    cols = [
        ("🏢  스튜디오", SAGE, [
            ("위치 데이터",     "lat / lon  ·  PostGIS ST_DWithin\nHaversine 거리 계산"),
            ("전문화 태그",     "prenatal / therapy / back care\nvinyasa / iyengar / senior"),
            ("운영 정보",       "영업시간 · 클래스 정원 · 가격"),
            ("평점",            "rating 4.0–5.0 (tiebreaker)"),
        ]),
        ("👩‍🏫  강사", EARTH, [
            ("자격증",          "RYT-200 / RYT-500\nPrenatal Certified\nYoga Therapist\nIyengar Certified"),
            ("계보 & 출처",     "lineage_source\nFYT100 세션 레퍼런스"),
            ("Schema.org",      "Person JSON-LD 자동 생성\nE-E-A-T 신호 → AI 검색 노출"),
            ("큐잉 스크립트",   "instructor_cue_priority\n부상별 수정 큐"),
        ]),
        ("🎬  콘텐츠", LOTUS, [
            ("유튜브 / 쇼츠",   "채널 태그 · 재생목록 slug\n조회수 신뢰도 가중"),
            ("인스타그램",      "릴스 해시태그\n강사 계정 연결"),
            ("GEO 북",          "OCR 텍스트 → 키워드 인덱스\nTesseract + cv2 파이프라인"),
            ("AEO JSON-LD",     "Schema.org Co-Author\nAI 검색 인용 최적화"),
        ]),
    ]

    for ci, (col_title, accent, rows) in enumerate(cols):
        cx = Inches(0.4 + ci * 4.3)
        cw = Inches(4.1)
        # column header
        rect(sl, cx, Inches(1.55), cw, Inches(0.48), HEADER_BG, accent, 1)
        txbox(sl, col_title, cx, Inches(1.55), cw, Inches(0.48),
              size=15, bold=True, color=accent, align=PP_ALIGN.CENTER)
        # rows
        for ri, (label, detail) in enumerate(rows):
            ry = Inches(2.1 + ri * 1.2)
            card(sl, cx, ry, cw, Inches(1.1), BG_CARD, accent, 1)
            txbox(sl, label, cx + Inches(0.12), ry + Inches(0.08), cw - Inches(0.2), Inches(0.38),
                  size=13, bold=True, color=accent)
            txbox(sl, detail, cx + Inches(0.12), ry + Inches(0.44), cw - Inches(0.2), Inches(0.6),
                  size=11, color=TEXT_DIM, wrap=True)

    # bottom note
    rect(sl, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.42), BG_CARD2, RULE_COLOR, 1)
    txbox(sl,
          "▶  프로바이더 온보딩: 스튜디오 등록 폼 → 운영자 검증 → DB 인입 → Kill-Switch 적용 → 매칭 엔진 반영",
          Inches(0.6), Inches(6.92), Inches(12.2), Inches(0.38),
          size=12, color=GOLD, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 사용자 여정: 4가지 발견 (User Journey)
# ═══════════════════════════════════════════════════════════════════════════════

def slide_user_journey():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, LOTUS)
    title_bar(sl, "04 · 사용자 여정 — 하나의 앱에서 4가지 발견", title_color=LOTUS)

    journeys = [
        (
            "01",
            "강사 찾기",
            "🧑‍🏫",
            SAGE,
            [
                "조건 입력: 허리 디스크 + 빌라넬 배경",
                "NeedFit × Cert Overlap 점수화",
                "강사 프로필: RYT-500 · Yoga Therapist",
                "Schema.org Person JSON-LD 임베딩",
                "→  Google SGE / Perplexity 노출",
            ],
        ),
        (
            "02",
            "유튜브 · 쇼츠 찾기",
            "📱",
            EARTH,
            [
                "증상 키워드 → GEO 북 RAG 검색",
                "채널 태그 · 재생목록 slug 매핑",
                "Elbee 챗봇 큐잉 해설 제공",
                "인스타 릴스 해시태그 연결",
                "→  영상 품질 가중치 (조회 신뢰도)",
            ],
        ),
        (
            "03",
            "스튜디오 찾기",
            "🏢",
            LOTUS,
            [
                "신체적 필요 × 위치 × 이동 거리 입력",
                "NeedFit 40% + Proximity 30%",
                "+ Specialization 30%",
                "Haversine 거리 계산 (PostGIS)",
                "→  지도 + 점수 투명 표시",
            ],
        ),
        (
            "04",
            "홈 큐잉 스크립트 + HUD",
            "🎙️",
            GOLD,
            [
                "안전 필터 통과한 포즈 시퀀스 선택",
                "instructor_cue_priority 기반 큐 생성",
                "Elbee → 한국어 해요체 보이스오버",
                "HUD 자막 오버레이 (AR 대응)",
                "→  집에서 강사 없이도 안전 수련",
            ],
        ),
    ]

    for i, (num, title, icon, accent, bullets) in enumerate(journeys):
        x = Inches(0.4 + i * 3.22)
        w = Inches(3.0)
        # number badge
        rect(sl, x, Inches(1.55), Inches(0.48), Inches(0.48), accent, line_color=None)
        txbox(sl, num, x, Inches(1.55), Inches(0.48), Inches(0.48),
              size=14, bold=True, color=BG_DEEP, align=PP_ALIGN.CENTER)
        # icon + title
        txbox(sl, icon, x + Inches(0.55), Inches(1.52), w - Inches(0.55), Inches(0.52),
              size=22, color=accent)
        txbox(sl, title, x + Inches(0.55), Inches(1.96), w - Inches(0.55), Inches(0.48),
              size=14, bold=True, color=accent)
        # bullet card
        card(sl, x, Inches(2.5), w, Inches(4.3), BG_CARD, accent, 1)
        y_b = Inches(2.65)
        for b in bullets:
            txbox(sl, f"• {b}", x + Inches(0.12), y_b, w - Inches(0.22), Inches(0.58),
                  size=12, color=TEXT_MAIN, wrap=True)
            y_b += Inches(0.72)

    # bottom CTA
    rect(sl, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.38), BG_CARD2, GOLD, 1)
    txbox(sl, "모든 여정은 단 하나의 조건 체크에서 시작됩니다 →  check.yogaman.club",
          Inches(0.6), Inches(6.97), Inches(12.0), Inches(0.35),
          size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 라이브 데모 흐름 (Demo Walk-through)
# ═══════════════════════════════════════════════════════════════════════════════

def slide_demo():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, GOLD)
    title_bar(sl, "05 · 라이브 데모 — 5단계 흐름", title_color=GOLD)

    steps = [
        ("STEP 1", "조건 입력",
         "Physical need 선택\n(허리 통증 / 산전 / 노인 등)\n위치 + 반경 설정",
         SAGE),
        ("STEP 2", "Kill-Switch 작동",
         "입력 조건과 일치하는\nkill_switch=true 포즈 즉시 제외\n점수 계산 전 실행",
         KILL_RED),
        ("STEP 3", "스튜디오 매칭",
         "NeedFit × Proximity\n× Specialization\n투명 점수 표시",
         EARTH),
        ("STEP 4", "강사 & 콘텐츠",
         "강사 자격증 카드\nSchema.org JSON-LD\n유튜브 / 쇼츠 링크",
         LOTUS),
        ("STEP 5", "홈 큐잉 생성",
         "Elbee 챗봇\n한국어 큐잉 스크립트\n보이스오버 HUD",
         GOLD),
    ]

    for i, (num, title, body, accent) in enumerate(steps):
        x = Inches(0.5 + i * 2.57)
        w = Inches(2.4)
        # step badge
        rect(sl, x, Inches(1.55), w, Inches(0.4), accent, line_color=None)
        txbox(sl, num, x, Inches(1.55), w, Inches(0.4),
              size=12, bold=True, color=BG_DEEP, align=PP_ALIGN.CENTER)
        card(sl, x, Inches(1.95), w, Inches(3.2), BG_CARD, accent, 1)
        txbox(sl, title, x + Inches(0.12), Inches(2.08), w - Inches(0.22), Inches(0.48),
              size=15, bold=True, color=accent)
        txbox(sl, body, x + Inches(0.12), Inches(2.58), w - Inches(0.22), Inches(2.4),
              size=13, color=TEXT_MAIN, wrap=True)

    # Live URL
    rect(sl, Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.58), BG_CARD2, SAGE, 1)
    txbox(sl, "🌐  match.yogaman.club  ·  elbee.yogaman.club  — 지금 바로 사용 가능",
          Inches(0.7), Inches(5.37), Inches(11.9), Inches(0.52),
          size=16, bold=True, color=SAGE, align=PP_ALIGN.CENTER)

    # Differentiation table
    txbox(sl, "경쟁 우위 요약", Inches(0.5), Inches(6.08), Inches(12.3), Inches(0.38),
          size=14, bold=True, color=TEXT_DIM)

    headers = ["기능", "유튜브/인스타", "Mindbody", "요가큐"]
    rows = [
        ["안전 Kill-Switch",     "✗", "✗", "✅"],
        ["강사 자격 검증",       "✗", "△", "✅  Schema.org"],
        ["홈 큐잉 보이스오버",   "✗", "✗", "✅  Elbee HUD"],
        ["AI 검색 노출 (AEO)",   "✗", "✗", "✅  JSON-LD"],
    ]
    col_widths = [Inches(4.5), Inches(2.2), Inches(2.2), Inches(3.2)]
    nrows, ncols = len(rows), len(headers)
    row_h = Inches(0.32)
    left = Inches(0.5)
    top = Inches(6.46)
    tbl = sl.shapes.add_table(nrows + 1, ncols, left, top,
                               sum(col_widths), row_h * (nrows + 1)).table
    for ci, w in enumerate(col_widths):
        tbl.columns[ci].width = w
    colors_h = [HEADER_BG] * ncols
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = HEADER_BG
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True; run.font.size = Pt(11)
        run.font.color.rgb = SAGE
    for ri, row_data in enumerate(rows):
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD if ri % 2 == 0 else BG_CARD2
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(11)
            color = SAGE if "✅" in val else (KILL_RED if "✗" in val else TEXT_DIM)
            run.font.color.rgb = color


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 기술 아키텍처 (Architecture visual)
# ═══════════════════════════════════════════════════════════════════════════════

def slide_architecture():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, SAGE)
    title_bar(sl, "06 · 기술 아키텍처 — 3개 파이프라인")

    lanes = [
        ("Pipeline A · 포즈 매칭", SAGE,
         ["yoga repo\n(poses_database)", "enrich_poses.py\nbenefit tags · kill_switch", "PostgreSQL\npose_contraindications", "Spring Boot\n매칭 엔진", "Top-K Poses\nmatch.yogaman.club"]),
        ("Pipeline B · 스튜디오 매칭", EARTH,
         ["User Input\nneed · lat/lon · km", "NeedFit 40%\n+ Proximity 30%", "+ Specialization 30%", "Studio Ranker\n가중합산", "Ranked Studios\nStreamlit Demo"]),
        ("Pipeline C · RAG 챗봇", LOTUS,
         ["screenshots/\n(book images)", "ocr_pipeline.py\nTesseract + cv2", "GeoDataStore\nkeyword index", "Elbee RAG\nOllama Mistral", "Streamed Answer\nelbee.yogaman.club"]),
    ]

    for li, (lane_title, accent, nodes) in enumerate(lanes):
        ly = Inches(1.55 + li * 1.8)
        # lane label
        rect(sl, Inches(0.4), ly, Inches(2.1), Inches(1.5), BG_CARD2, accent, 1)
        txbox(sl, lane_title, Inches(0.4), ly + Inches(0.5), Inches(2.1), Inches(0.5),
              size=11, bold=True, color=accent, align=PP_ALIGN.CENTER, wrap=True)

        # nodes
        nw = Inches(1.98)
        for ni, node in enumerate(nodes):
            nx = Inches(2.6 + ni * 2.12)
            card(sl, nx, ly, nw, Inches(1.5), BG_CARD, accent, 1)
            txbox(sl, node, nx + Inches(0.1), ly + Inches(0.35), nw - Inches(0.18), Inches(0.75),
                  size=11, color=TEXT_MAIN, align=PP_ALIGN.CENTER, wrap=True)
            # arrow
            if ni < len(nodes) - 1:
                ax = nx + nw + Inches(0.02)
                rect(sl, ax, ly + Inches(0.73), Inches(0.1), Pt(2.5), accent)

    # Kill-Switch callout
    rect(sl, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.34), BG_CARD2, KILL_RED, 1)
    txbox(sl, "⛔  Kill-Switch 불변식: kill_switch=true 포즈는 점수 계산 이전에 완전 제외 — 어떤 점수도 이를 override 할 수 없습니다",
          Inches(0.6), Inches(7.07), Inches(12.1), Inches(0.3),
          size=11, bold=True, color=KILL_RED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 기술 해자 (Moat)
# ═══════════════════════════════════════════════════════════════════════════════

def slide_moat():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, GOLD)
    title_bar(sl, "07 · 기술 해자 — 복제하기 어려운 이유", title_color=GOLD)

    moats = [
        ("⛔  Safety-First\nKill-Switch",
         "CRITICAL 등급 금기 포즈를 점수 전에 제거.\n"
         "이 논리를 잘못 구현하면 법적 책임 → 진입 장벽.",
         KILL_RED),
        ("🏷️  E-E-A-T\n강사 데이터",
         "Schema.org Person JSON-LD 자동 생성.\n"
         "Google SGE · Perplexity가 인용 = 강사 신뢰 자산 누적.",
         SAGE),
        ("🗺️  Haversine +\n전문화 복합 매칭",
         "거리 × 자격증 × 전문화 × 평점을 동시 계산.\n"
         "단순 리스트뷰로는 재현 불가.",
         EARTH),
        ("🎙️  HUD 큐잉\n보이스오버",
         "Elbee가 해요체 큐잉 스크립트를 실시간 생성.\n"
         "Ollama 실패 → OpenAI 자동 전환 (99.9% uptime).",
         LOTUS),
        ("📖  GEO 콘텐츠\n지식 그래프",
         "OCR → 키워드 인덱스 → RAG → AEO 최적화.\n"
         "AI 검색 인용 구조는 수년간 쌓인 콘텐츠 자산.",
         GOLD),
    ]

    for i, (title, body, accent) in enumerate(moats):
        x = Inches(0.4 + (i % 3) * 4.3)
        y = Inches(1.55 + (i // 3) * 2.5)
        card(sl, x, y, Inches(4.0), Inches(2.2), BG_CARD, accent, 1)
        txbox(sl, title, x + Inches(0.15), y + Inches(0.15), Inches(3.7), Inches(0.72),
              size=15, bold=True, color=accent, wrap=True)
        txbox(sl, body, x + Inches(0.15), y + Inches(0.85), Inches(3.7), Inches(1.2),
              size=12, color=TEXT_DIM, wrap=True)

    txbox(sl, "Network effect: 강사·스튜디오 데이터가 많을수록 매칭 정확도 상승 → 신규 진입자 불리",
          Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4),
          size=13, italic=True, color=GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — 비즈니스 모델 & 트랙션
# ═══════════════════════════════════════════════════════════════════════════════

def slide_biz():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, EARTH)
    title_bar(sl, "08 · 비즈니스 모델 & 현재 트랙션", title_color=EARTH)

    # Revenue streams
    txbox(sl, "매출 라인", Inches(0.55), Inches(1.55), Inches(5.8), Inches(0.45),
          size=16, bold=True, color=EARTH)
    revenue_rows = [
        ("B2B SaaS",     "스튜디오 구독  ₩49,000/월  (리드 + 예약 API)"),
        ("B2C 프리미엄", "사용자 맞춤 큐잉 · HUD 보이스오버  ₩9,900/월"),
        ("API 라이선스", "Mindbody · 1club 파트너 매칭 API"),
        ("AEO 컨설팅",   "강사 Schema.org 페이지 구축 대행"),
    ]
    y = Inches(2.05)
    for rev_type, detail in revenue_rows:
        rect(sl, Inches(0.55), y, Inches(1.7), Inches(0.38), HEADER_BG, EARTH, 1)
        txbox(sl, rev_type, Inches(0.55), y, Inches(1.7), Inches(0.38),
              size=12, bold=True, color=EARTH, align=PP_ALIGN.CENTER)
        txbox(sl, detail, Inches(2.3), y, Inches(4.0), Inches(0.38),
              size=12, color=TEXT_MAIN)
        y += Inches(0.48)

    # Traction stats
    txbox(sl, "현재 운영 지표", Inches(6.9), Inches(1.55), Inches(6.0), Inches(0.45),
          size=16, bold=True, color=SAGE)
    stats = [
        ("856+",    "포즈 데이터\n(E-E-A-T 스키마)"),
        ("8",       "서울 스튜디오\n(MVP 파트너)"),
        ("2",       "라이브 도메인\n(Cloudflare)"),
        ("99.9%",   "챗봇 가용성\n(Ollama+OpenAI)"),
    ]
    for i, (num, label) in enumerate(stats):
        sx = Inches(7.0 + i * 1.55)
        rect(sl, sx, Inches(2.1), Inches(1.35), Inches(1.6), BG_CARD, SAGE, 1)
        txbox(sl, num, sx, Inches(2.2), Inches(1.35), Inches(0.65),
              size=26, bold=True, color=SAGE, align=PP_ALIGN.CENTER)
        txbox(sl, label, sx, Inches(2.85), Inches(1.35), Inches(0.75),
              size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER, wrap=True)

    # GTM
    txbox(sl, "시장 진입 전략", Inches(0.55), Inches(4.3), Inches(12.3), Inches(0.45),
          size=15, bold=True, color=TEXT_DIM)
    gtm = [
        ("01 스튜디오 확보", "요가 스튜디오 파트너십 → 리드 데이터 수집", SAGE),
        ("02 쿼리 수확",     "AI 검색 트래픽 → GEO 콘텐츠 최적화 → 사용자 유입", EARTH),
        ("03 소비자 성장",   "홈 큐잉 HUD → B2C 구독 전환", LOTUS),
    ]
    for i, (step, desc, accent) in enumerate(gtm):
        x = Inches(0.55 + i * 4.25)
        card(sl, x, Inches(4.85), Inches(4.0), Inches(1.4), BG_CARD, accent, 1)
        txbox(sl, step, x + Inches(0.12), Inches(4.98), Inches(3.75), Inches(0.42),
              size=13, bold=True, color=accent)
        txbox(sl, desc, x + Inches(0.12), Inches(5.4), Inches(3.75), Inches(0.7),
              size=12, color=TEXT_DIM, wrap=True)

    # Ask
    rect(sl, Inches(0.55), Inches(6.45), Inches(12.3), Inches(0.78), BG_CARD2, GOLD, 1)
    txbox(sl, "💰  프리시드 라운드  ·  모집 금액: ₩3억  ·  용도: 강사 DB 구축 · HUD 보이스오버 개발 · GTM",
          Inches(0.75), Inches(6.58), Inches(11.9), Inches(0.55),
          size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Closing
# ═══════════════════════════════════════════════════════════════════════════════

def slide_closing():
    sl = prs.slides.add_slide(BLANK)
    bg(sl)
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, SAGE)
    # large visual quote
    rect(sl, Inches(0.4), Inches(1.2), Inches(12.5), Inches(2.8), BG_CARD, SAGE, 1)
    txbox(sl,
          '"안전을 먼저, 효과를 그 다음에."',
          Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.1),
          size=36, bold=True, italic=True, color=SAGE, align=PP_ALIGN.CENTER)
    txbox(sl,
          "요가큐는 몸 상태를 이해하고, 안전한 포즈·강사·스튜디오·홈 큐잉을 연결하는 유일한 AI 매칭 플랫폼입니다.",
          Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.1),
          size=18, color=TEXT_DIM, align=PP_ALIGN.CENTER, wrap=True)

    # links
    for i, (label, url) in enumerate([
        ("매칭 API", "match.yogaman.club"),
        ("AI 챗봇", "elbee.yogaman.club"),
        ("GitHub", "github.com/aiegoo/aeogeo"),
    ]):
        rect(sl, Inches(2.5 + i * 2.9), Inches(4.3), Inches(2.6), Inches(0.65),
             BG_CARD2, SAGE, 1)
        txbox(sl, label, Inches(2.5 + i * 2.9), Inches(4.35), Inches(2.6), Inches(0.28),
              size=11, bold=True, color=SAGE, align=PP_ALIGN.CENTER)
        txbox(sl, url, Inches(2.5 + i * 2.9), Inches(4.63), Inches(2.6), Inches(0.28),
              size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    txbox(sl, "함께 AI 요가 큐레이터를 만들어요.",
          Inches(0.4), Inches(5.3), Inches(12.5), Inches(0.8),
          size=28, bold=True, color=TEXT_MAIN, align=PP_ALIGN.CENTER)
    txbox(sl, "hello@yogaman.club  ·  AIBE5 · Team 14",
          Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.45),
          size=15, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # decorative leaves line at bottom
    rect(sl, Inches(3.5), Inches(6.9), Inches(6.3), Pt(2.5), SAGE)


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════════════

slide_cover()
slide_problem()
slide_condition_check()
slide_provider_data()
slide_user_journey()
slide_demo()
slide_architecture()
slide_moat()
slide_biz()
slide_closing()

OUT = "/mnt/c/Users/hsyyu/Downloads/요가큐_신규발표.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
print(f"Slides: {len(prs.slides)}")
