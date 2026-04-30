#!/usr/bin/env python3
"""
combine_pptx.py — Merge two PPTX files.
Order: aeogeo_pitch_deck_ko.pptx (v2, 13 slides) → 요가큐_발표.pptx (v1, 12 slides)
Output: /mnt/c/Users/hsyyu/Downloads/요가큐_발표_combined.pptx
"""

from pptx import Presentation
from pptx.util import Emu
from lxml import etree
from io import BytesIO
import copy

P_NS  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
A_NS  = 'http://schemas.openxmlformats.org/drawingml/2006/main'
R_NS  = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

V2_PATH  = '/mnt/c/Users/hsyyu/Downloads/aeogeo_pitch_deck_ko.pptx'
V1_PATH  = '/mnt/c/Users/hsyyu/Downloads/요가큐_발표.pptx'
OUT_PATH = '/mnt/c/Users/hsyyu/Downloads/요가큐_발표_combined.pptx'


def copy_slide(src_slide, dest_prs):
    """Append a deep-copy of src_slide into dest_prs, including embedded images."""
    blank = dest_prs.slide_layouts[6]          # blank layout
    dest_slide = dest_prs.slides.add_slide(blank)

    # ── 1. Copy image/media relationships first, build old→new rId map ────────
    rId_map = {}
    for rId, rel in list(src_slide.part.rels.items()):
        if '/image' in rel.reltype:
            try:
                blob = rel.target_part.blob
                _, new_rId = dest_slide.part.get_or_add_image_part(BytesIO(blob))
                rId_map[rId] = new_rId
            except Exception:
                pass  # skip rels we can't transfer

    # ── 2. Replace spTree content, remapping r:embed to new rIds ──────────────
    R_EMBED = f'{{{R_NS}}}embed'
    R_LINK  = f'{{{R_NS}}}link'
    dest_spTree = dest_slide.shapes._spTree
    for child in list(dest_spTree):
        dest_spTree.remove(child)
    for child in src_slide.shapes._spTree:
        node = copy.deepcopy(child)
        for el in node.iter():
            for attr in (R_EMBED, R_LINK):
                val = el.get(attr)
                if val and val in rId_map:
                    el.set(attr, rId_map[val])
        dest_spTree.append(node)

    # ── 3. Copy slide background (p:bg inside p:cSld) ─────────────────────────
    src_cSld  = src_slide._element.find(f'{{{P_NS}}}cSld')
    dest_cSld = dest_slide._element.find(f'{{{P_NS}}}cSld')
    if src_cSld is not None and dest_cSld is not None:
        src_bg = src_cSld.find(f'{{{P_NS}}}bg')
        if src_bg is not None:
            dest_bg = dest_cSld.find(f'{{{P_NS}}}bg')
            if dest_bg is not None:
                dest_cSld.remove(dest_bg)
            dest_cSld.insert(0, copy.deepcopy(src_bg))

    return dest_slide


# ── Load sources ──────────────────────────────────────────────────────────────
v2 = Presentation(V2_PATH)
v1 = Presentation(V1_PATH)

# ── Build combined (use v2 dimensions as base) ────────────────────────────────
combined = Presentation()
combined.slide_width  = v2.slide_width
combined.slide_height = v2.slide_height

# ── Append v2 slides first ────────────────────────────────────────────────────
for i, slide in enumerate(v2.slides):
    copy_slide(slide, combined)
    print(f"  v2 [{i+1:02d}/{len(v2.slides)}] copied")

# ── Separator slide ───────────────────────────────────────────────────────────
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

sep = combined.slides.add_slide(combined.slide_layouts[6])
fill = sep.background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

sep_box = sep.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.5))
sep_tf  = sep_box.text_frame
sep_tf.word_wrap = True
sep_p   = sep_tf.paragraphs[0]
sep_p.alignment = PP_ALIGN.CENTER
sep_run = sep_p.add_run()
sep_run.text = "발표 스크립트 상세"
sep_run.font.size  = Pt(48)
sep_run.font.bold  = True
sep_run.font.color.rgb = RGBColor(0xa7, 0x8b, 0xfa)

sep_sub = sep.shapes.add_textbox(Inches(1), Inches(4.1), Inches(11.3), Inches(0.7))
sep_sub.text_frame.word_wrap = True
sub_p = sep_sub.text_frame.paragraphs[0]
sub_p.alignment = PP_ALIGN.CENTER
sub_run = sub_p.add_run()
sub_run.text = "요가큐 — 기술·데이터·운영 발표 자료"
sub_run.font.size  = Pt(24)
sub_run.font.color.rgb = RGBColor(0x7d, 0xd3, 0xfc)
print(f"  separator slide inserted")

# ── Append v1 slides after ────────────────────────────────────────────────────
for i, slide in enumerate(v1.slides):
    copy_slide(slide, combined)
    print(f"  v1 [{i+1:02d}/{len(v1.slides)}] copied")

# ── Save ──────────────────────────────────────────────────────────────────────
combined.save(OUT_PATH)
total = len(v2.slides) + 1 + len(v1.slides)
print(f"\n✅  Saved → {OUT_PATH}")
print(f"   {len(v2.slides)} (v2)  +  1 (separator)  +  {len(v1.slides)} (v1)  =  {total} slides total")
