#!/usr/bin/env python3
"""
build_yoga_deck_v2.py
─────────────────────
Reference theme: light pastel (DevBridge-style) — blue/cream rounded pills.
+ matplotlib architecture diagram embedded
+ real demo screenshots embedded
+ 10 slides, full differentiator content
Output: /mnt/c/Users/hsyyu/Downloads/요가큐_발표_v2.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Reference-style Pastel Palette ───────────────────────────────────────────
BG_LIGHT   = RGBColor(0xEE, 0xF4, 0xFB)   # light sky blue bg
BG_GRAD1   = RGBColor(0xD8, 0xE8, 0xF8)   # top gradient band
BG_GRAD2   = RGBColor(0xE8, 0xDC, 0xF5)   # bottom gradient band (lavender)
CARD_WHITE = RGBColor(0xFF, 0xFF, 0xFF)   # white card
CARD_CREAM = RGBColor(0xFB, 0xFB, 0xEE)   # cream solution box
CARD_BLUE  = RGBColor(0xE8, 0xF2, 0xFD)   # soft blue card
CARD_DARK  = RGBColor(0x1E, 0x2D, 0x45)   # dark navy (accent section)

BLUE_TITLE = RGBColor(0x4A, 0x7E, 0xBB)   # steel blue title
BLUE_PILL  = RGBColor(0x6B, 0x9E, 0xD4)   # pill label blue
BLUE_RULE  = RGBColor(0x9B, 0xC4, 0xE8)   # rule / border
BLUE_DARK  = RGBColor(0x2C, 0x5F, 0x9E)   # dark blue accent

SAGE       = RGBColor(0x5A, 0x9A, 0x6E)   # green (safety / ok)
RED_KILL   = RGBColor(0xC0, 0x3C, 0x3C)   # kill-switch red
EARTH      = RGBColor(0xB5, 0x6A, 0x35)   # terracotta
GOLD       = RGBColor(0xB8, 0x8C, 0x2A)   # warm gold
LOTUS      = RGBColor(0x7B, 0x5E, 0xAA)   # purple

TEXT_DARK  = RGBColor(0x1A, 0x2A, 0x3A)   # near-black text
TEXT_MID   = RGBColor(0x4A, 0x5A, 0x6A)   # mid-gray text
TEXT_DIM   = RGBColor(0x7A, 0x8A, 0x9A)   # dim gray

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

ARCH_PNG    = '/home/aiegoo/repos/aiegoo/aeogeo/assets/architecture_diagram.png'
SHOT_RANKED = '/home/aiegoo/repos/aiegoo/aeogeo/assets/teaser/shot_04_ranked.png'
SHOT_SCORE  = '/home/aiegoo/repos/aiegoo/aeogeo/assets/teaser/shot_05_breakdown.png'
SHOT_JSONLD = '/home/aiegoo/repos/aiegoo/aeogeo/assets/teaser/shot_07_jsonld.png'


# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color=BG_LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill_color, line_color=None, line_pt=0.75):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s


def pill(slide, left, top, width, height, fill=CARD_BLUE, border=BLUE_PILL, border_pt=1.2):
    """Rounded rectangle approximated via MSO_SHAPE_TYPE rounded_rectangle (id=5)"""
    s = slide.shapes.add_shape(5, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(border_pt)
    # adjust corner radius
    adj = s.adjustments
    if adj:
        adj[0] = 0.08  # roundness
    return s


def txbox(slide, text, left, top, width, height,
          size=14, bold=False, italic=False, color=TEXT_DARK,
          align=PP_ALIGN.LEFT, wrap=True):
    b = slide.shapes.add_textbox(left, top, width, height)
    tf = b.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return b


def multiline(slide, lines, left, top, width, height, wrap=True):
    """lines: list of (text, size, bold, color, align)"""
    b = slide.shapes.add_textbox(left, top, width, height)
    tf = b.text_frame; tf.word_wrap = wrap
    first = True
    for (text, size, bold, color, align) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return b


def title_bar(slide, title, subtitle=None, title_color=BLUE_TITLE):
    # top rule band
    rect(slide, 0, 0, SLIDE_W, Inches(0.12), BG_GRAD1)
    txbox(slide, title,
          Inches(0.55), Inches(0.2), Inches(12.0), Inches(0.72),
          size=28, bold=True, color=title_color)
    rect(slide, Inches(0.55), Inches(0.92), Inches(12.2), Pt(2), BLUE_RULE)
    if subtitle:
        txbox(slide, subtitle,
              Inches(0.55), Inches(0.97), Inches(12.2), Inches(0.4),
              size=13, color=TEXT_MID, italic=True)


def problem_pill(slide, number, problem, solution, y):
    """Renders one Problem → Solution row like the reference."""
    # problem pill
    pill(slide, Inches(0.55), y, Inches(6.2), Inches(0.82),
         fill=CARD_BLUE, border=BLUE_PILL, border_pt=1.2)
    txbox(slide, number,
          Inches(0.62), y + Inches(0.18), Inches(1.1), Inches(0.45),
          size=13, bold=True, color=BLUE_PILL, align=PP_ALIGN.CENTER)
    # vertical divider
    rect(slide, Inches(1.8), y + Inches(0.1), Pt(1.5), Inches(0.62), BLUE_RULE)
    txbox(slide, problem,
          Inches(1.92), y + Inches(0.1), Inches(4.7), Inches(0.65),
          size=13, color=TEXT_DARK, wrap=True)
    # arrow >>
    txbox(slide, '>>',
          Inches(6.95), y + Inches(0.22), Inches(0.5), Inches(0.4),
          size=18, bold=True, color=BLUE_PILL, align=PP_ALIGN.CENTER)
    # solution pill (cream)
    pill(slide, Inches(7.6), y, Inches(5.15), Inches(0.82),
         fill=CARD_CREAM, border=BLUE_PILL, border_pt=1.0)
    txbox(slide, solution,
          Inches(7.72), y + Inches(0.18), Inches(4.9), Inches(0.45),
          size=13, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)


def step_box(slide, step_label, title, body, left, top, width=Inches(2.35), accent=BLUE_DARK):
    rect(slide, left, top, width, Inches(0.38), accent)
    txbox(slide, step_label, left, top, width, Inches(0.38),
          size=11, bold=True, color=CARD_WHITE, align=PP_ALIGN.CENTER)
    pill(slide, left, top + Inches(0.38), width, Inches(2.6),
         fill=CARD_WHITE, border=accent, border_pt=1.2)
    txbox(slide, title, left + Inches(0.1), top + Inches(0.52),
          width - Inches(0.18), Inches(0.45),
          size=13, bold=True, color=accent)
    txbox(slide, body, left + Inches(0.1), top + Inches(0.97),
          width - Inches(0.18), Inches(1.85),
          size=11, color=TEXT_MID, wrap=True)


def img(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        return
    if width and height:
        slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(path, left, top, height=height)
    else:
        slide.shapes.add_picture(path, left, top)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════════════════════════════════

def slide_cover():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, BG_LIGHT)
    # gradient bands top & bottom
    rect(sl, 0, 0, SLIDE_W, Inches(2.2), BG_GRAD1)
    rect(sl, 0, Inches(5.6), SLIDE_W, Inches(1.9), BG_GRAD2)
    # left accent stripe
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, BLUE_DARK)

    txbox(sl, "요가큐 · Yogaman.club",
          Inches(0.5), Inches(0.55), Inches(9), Inches(1.0),
          size=42, bold=True, color=BLUE_DARK)
    txbox(sl, "AI 기반 요가 안전 매칭 플랫폼",
          Inches(0.5), Inches(1.55), Inches(9), Inches(0.6),
          size=22, color=BLUE_TITLE)
    txbox(sl, "몸 상태를 이해하고 안전한 포즈 · 강사 · 스튜디오 · 홈 큐잉을 연결합니다.",
          Inches(0.5), Inches(2.25), Inches(9), Inches(0.65),
          size=16, color=TEXT_MID, wrap=True)

    # 5-step flow card on right
    pill(sl, Inches(9.8), Inches(0.55), Inches(3.3), Inches(5.5),
         fill=CARD_WHITE, border=BLUE_RULE, border_pt=1.5)
    txbox(sl, "5단계 흐름", Inches(9.9), Inches(0.72), Inches(3.1), Inches(0.42),
          size=14, bold=True, color=BLUE_TITLE, align=PP_ALIGN.CENTER)
    steps = [
        ("① 몸 상태 체크",        SAGE),
        ("② Kill-Switch 안전 필터", RED_KILL),
        ("③ 포즈 / 스튜디오 매칭", BLUE_DARK),
        ("④ 강사 · 쇼츠 발견",     EARTH),
        ("⑤ 홈 큐잉 보이스오버",   LOTUS),
    ]
    for i, (t, c) in enumerate(steps):
        pill(sl, Inches(10.0), Inches(1.2 + i * 0.85), Inches(2.9), Inches(0.62),
             fill=CARD_BLUE, border=BLUE_RULE, border_pt=0.8)
        txbox(sl, t, Inches(10.1), Inches(1.3 + i * 0.85), Inches(2.7), Inches(0.42),
              size=13, bold=True, color=c, align=PP_ALIGN.CENTER)

    # bottom row
    for i, (label, url) in enumerate([
        ("매칭 API", "match.yogaman.club"),
        ("AI 챗봇",  "elbee.yogaman.club"),
    ]):
        pill(sl, Inches(0.5 + i * 4.0), Inches(6.35), Inches(3.5), Inches(0.65),
             fill=CARD_CREAM, border=BLUE_PILL, border_pt=1.0)
        txbox(sl, f"{label}  ·  {url}",
              Inches(0.6 + i * 4.0), Inches(6.48), Inches(3.3), Inches(0.38),
              size=13, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
    txbox(sl, "AIBE5 · Team 14  ·  April 2026",
          Inches(9.0), Inches(7.1), Inches(4.2), Inches(0.32),
          size=11, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem / Solution (DevBridge-style pill layout)
# ═══════════════════════════════════════════════════════════════════════════════

def slide_problem():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, BG_LIGHT)
    rect(sl, 0, 0, SLIDE_W, Inches(0.12), BG_GRAD1)
    rect(sl, 0, Inches(7.38), SLIDE_W, Inches(0.12), BG_GRAD2)

    txbox(sl, "핵심 문제와 솔루션",
          Inches(0.55), Inches(0.22), Inches(8), Inches(0.65),
          size=28, bold=True, color=BLUE_TITLE)
    txbox(sl, "Solution",
          Inches(7.9), Inches(0.28), Inches(5.0), Inches(0.5),
          size=18, color=BLUE_PILL, align=PP_ALIGN.CENTER)
    rect(sl, Inches(0.55), Inches(0.88), Inches(12.2), Pt(1.8), BLUE_RULE)

    problems = [
        ("Problem 01",
         "'일감 매칭'에만 치중된\n기존 요가 앱들",
         "안전 우선 설계 — Kill-Switch 금기 필터"),
        ("Problem 02",
         "강사 자격 검증 없이\n알고리즘이 추천",
         "Schema.org E-E-A-T 강사 신뢰 자산"),
        ("Problem 03",
         "스튜디오 = 리스트뷰,\n복합 매칭 없음",
         "NeedFit × 거리 × 전문화 3요소 매칭"),
    ]
    for i, (num, prob, sol) in enumerate(problems):
        problem_pill(sl, num, prob, sol, Inches(1.2 + i * 1.85))

    # stats footer
    rect(sl, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.5), BG_GRAD1)
    stats = [("1,150억 $", "글로벌 요가 시장"), ("68%", "앱 이탈률"), ("0", "안전 필터 앱")]
    for i, (n, l) in enumerate(stats):
        txbox(sl, n, Inches(1.5 + i * 3.8), Inches(6.87), Inches(3.0), Inches(0.28),
              size=18, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
        txbox(sl, l, Inches(1.5 + i * 3.8), Inches(7.13), Inches(3.0), Inches(0.28),
              size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — 몸 상태 체크 → 메타데이터 매핑
# ═══════════════════════════════════════════════════════════════════════════════

def slide_condition_check():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, BG_LIGHT)
    rect(sl, 0, 0, SLIDE_W, Inches(0.12), BG_GRAD1)
    title_bar(sl, "02 · 몸 상태 체크 → 메타데이터 매핑 → 안전 후보군", title_color=BLUE_TITLE)

    flow = [
        ("STEP 1\n사용자 입력",  "• 허리 디스크\n• 고혈압\n• 임신 중\n• 무릎 수술 후",       BLUE_DARK),
        ("STEP 2\n컨디션 태그",  "disc / hypertension\npregnancy / knee_surgery\n→ DB 조회",     BLUE_TITLE),
        ("STEP 3\nKill-Switch",  "kill_switch = TRUE\n→ 완전 제외\n점수 계산 전 실행",          RED_KILL),
        ("STEP 4\n안전 후보군",  "CAUTION: 수정 큐잉\nMEDICAL_CLEARANCE:\n의사 확인 권고",     SAGE),
        ("STEP 5\n매칭 & 순위",  "benefit.weight\n× goal_overlap\n→ Top-K 반환",               LOTUS),
    ]
    bw = Inches(2.3)
    for i, (label, body, accent) in enumerate(flow):
        x = Inches(0.38 + i * 2.58)
        step_box(sl, label, "", body, x, Inches(1.35), bw, accent)
        if i < len(flow) - 1:
            txbox(sl, "→", x + bw + Inches(0.08), Inches(2.4),
                  Inches(0.38), Inches(0.5), size=20, bold=True,
                  color=BLUE_PILL, align=PP_ALIGN.CENTER)

    # schema strip
    rect(sl, Inches(0.38), Inches(5.3), Inches(12.55), Inches(1.15), CARD_CREAM, BLUE_RULE, 0.8)
    txbox(sl, "📐  메타데이터 스키마  (pose_eat_schema.json)",
          Inches(0.55), Inches(5.38), Inches(12.2), Inches(0.38),
          size=13, bold=True, color=BLUE_DARK)
    txbox(sl,
          'contraindications[]:  { condition,  severity: "CAUTION | CRITICAL | MEDICAL_CLEARANCE",  '
          'kill_switch: bool,  instruction }    ·    '
          'benefits[]:  { tag,  weight: 0–1 }    ·    '
          'difficulty_rank: 1–5    ·    Schema.org ExerciseAction JSON-LD 자동 생성',
          Inches(0.55), Inches(5.76), Inches(12.2), Inches(0.62),
          size=11, color=TEXT_MID, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — 프로바이더 데이터
# ═══════════════════════════════════════════════════════════════════════════════

def slide_provider_data():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, BG_LIGHT)
    rect(sl, 0, 0, SLIDE_W, Inches(0.12), BG_GRAD1)
    title_bar(sl, "03 · 프로바이더 데이터 — 매칭 엔진을 작동시키는 원천")

    cols = [
        ("🏢  스튜디오", BLUE_DARK, [
            ("위치 데이터",   "lat / lon · PostGIS\nHaversine 거리"),
            ("전문화 태그",   "prenatal / therapy\nback care / senior"),
            ("운영 정보",     "영업시간 · 정원 · 가격"),
            ("평점",          "4.0–5.0 (tiebreaker)"),
        ]),
        ("👩‍🏫  강사", EARTH, [
            ("자격증",        "RYT-200 / RYT-500\nPrenatal · Therapist"),
            ("계보 출처",     "lineage_source\nFYT100 레퍼런스"),
            ("Schema.org",    "Person JSON-LD\nE-E-A-T AI 검색 노출"),
            ("큐잉 스크립트", "instructor_cue_priority\n부상별 수정 큐"),
        ]),
        ("🎬  콘텐츠", LOTUS, [
            ("유튜브 / 쇼츠", "채널 태그 · 재생목록\n조회수 신뢰 가중"),
            ("인스타그램",    "릴스 해시태그\n강사 계정 연결"),
            ("GEO 북",        "OCR 텍스트 키워드 인덱스\nTesseract + cv2"),
            ("AEO JSON-LD",   "Schema.org Co-Author\nAI 인용 최적화"),
        ]),
    ]

    for ci, (col_title, accent, rows) in enumerate(cols):
        cx = Inches(0.38 + ci * 4.3)
        cw = Inches(4.1)
        rect(sl, cx, Inches(1.35), cw, Inches(0.48), accent)
        txbox(sl, col_title, cx, Inches(1.35), cw, Inches(0.48),
              size=15, bold=True, color=CARD_WHITE, align=PP_ALIGN.CENTER)
        for ri, (label, detail) in enumerate(rows):
            ry = Inches(1.92 + ri * 1.25)
            pill(sl, cx, ry, cw, Inches(1.12), CARD_WHITE, accent, 0.8)
            txbox(sl, label, cx + Inches(0.12), ry + Inches(0.1),
                  cw - Inches(0.22), Inches(0.38), size=13, bold=True, color=accent)
            txbox(sl, detail, cx + Inches(0.12), ry + Inches(0.48),
                  cw - Inches(0.22), Inches(0.56), size=11, color=TEXT_MID, wrap=True)

    rect(sl, Inches(0.38), Inches(7.0), Inches(12.55), Inches(0.38), BG_GRAD2)
    txbox(sl,
          "▶  프로바이더 온보딩: 등록 폼 → 운영자 검증 → DB 인입 → Kill-Switch 적용 → 매칭 반영",
          Inches(0.55), Inches(7.02), Inches(12.2), Inches(0.34),
          size=12, bold=True, color=BLUE_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 사용자 여정: 4가지 발견
# ═══════════════════════════════════════════════════════════════════════════════

def slide_user_journey():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, BG_LIGHT)
    rect(sl, 0, 0, SLIDE_W, Inches(0.12), BG_GRAD1)
    title_bar(sl, "04 · 사용자 여정 — 하나의 앱에서 4가지 발견")

    journeys = [
        ("01", "강사 찾기", "🧑‍🏫", BLUE_DARK,
         ["조건 입력: 허리 디스크 + 빌라넬",
          "NeedFit × Cert Overlap 점수화",
          "강사 프로필: RYT-500",
          "Schema.org Person JSON-LD",
          "→ Google SGE / Perplexity 노출"]),
        ("02", "유튜브·쇼츠 찾기", "📱", EARTH,
         ["증상 키워드 → GEO 북 RAG",
          "채널 태그 · 재생목록 매핑",
          "Elbee 챗봇 큐잉 해설",
          "인스타 릴스 해시태그 연결",
          "→ 영상 품질 가중 추천"]),
        ("03", "스튜디오 찾기", "🏢", LOTUS,
         ["필요 × 위치 × 이동 거리 입력",
          "NeedFit 40% + Proximity 30%",
          "+ Specialization 30%",
          "Haversine 거리 (PostGIS)",
          "→ 지도 + 점수 투명 표시"]),
        ("04", "홈 큐잉 스크립트 + HUD", "🎙️", SAGE,
         ["안전 필터 통과 포즈 시퀀스",
          "instructor_cue_priority 기반",
          "Elbee → 한국어 해요체 큐잉",
          "HUD 자막 오버레이 (AR 대응)",
          "→ 집에서 강사 없이 안전 수련"]),
    ]

    for i, (num, title, icon, accent, bullets) in enumerate(journeys):
        x = Inches(0.38 + i * 3.22)
        w = Inches(3.05)
        # header pill
        pill(sl, x, Inches(1.35), w, Inches(0.78), CARD_BLUE, accent, 1.2)
        txbox(sl, num, x + Inches(0.12), Inches(1.42), Inches(0.5), Inches(0.55),
              size=18, bold=True, color=accent)
        txbox(sl, f"{icon}  {title}", x + Inches(0.65), Inches(1.46),
              w - Inches(0.72), Inches(0.55),
              size=14, bold=True, color=accent)
        # bullet card
        pill(sl, x, Inches(2.22), w, Inches(4.6), CARD_WHITE, accent, 1.0)
        for j, b in enumerate(bullets):
            txbox(sl, f"• {b}", x + Inches(0.15), Inches(2.38 + j * 0.78),
                  w - Inches(0.25), Inches(0.7), size=12, color=TEXT_DARK, wrap=True)

    rect(sl, Inches(0.38), Inches(7.0), Inches(12.55), Inches(0.38), BG_GRAD2)
    txbox(sl, "모든 여정은 단 하나의 조건 체크에서 시작됩니다 →  check.yogaman.club",
          Inches(0.55), Inches(7.02), Inches(12.2), Inches(0.34),
          size=13, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 라이브 데모 + 실제 스크린샷
# ═══════════════════════════════════════════════════════════════════════════════

def slide_demo():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, BG_LIGHT)
    rect(sl, 0, 0, SLIDE_W, Inches(0.12), BG_GRAD1)
    title_bar(sl, "05 · 라이브 데모 — 실제 화면")

    # Left: 3-step flow
    steps = [
        ("조건 입력 & Kill-Switch",
         "Physical need 선택 → kill_switch=TRUE 포즈 즉시 제외\n점수 계산 전 실행. 안전 우선 설계.",
         RED_KILL),
        ("스튜디오 매칭 점수",
         "NeedFit 40% + Proximity 30% + Specialization 30%\n투명한 점수 표시. 8개 서울 스튜디오 MVP.",
         BLUE_DARK),
        ("Schema.org JSON-LD",
         "강사·포즈 페이지에 ExerciseAction/Person 자동 임베딩\nGoogle SGE · Perplexity AI 검색 인용.",
         LOTUS),
    ]
    for i, (title, body, accent) in enumerate(steps):
        pill(sl, Inches(0.38), Inches(1.38 + i * 1.78), Inches(5.6), Inches(1.58),
             CARD_WHITE, accent, 1.2)
        txbox(sl, title,
              Inches(0.55), Inches(1.52 + i * 1.78), Inches(5.25), Inches(0.44),
              size=14, bold=True, color=accent)
        txbox(sl, body,
              Inches(0.55), Inches(1.96 + i * 1.78), Inches(5.25), Inches(0.85),
              size=12, color=TEXT_MID, wrap=True)

    # Right: actual screenshots
    shot_y = Inches(1.35)
    for path, cap, w, h in [
        (SHOT_RANKED, "스튜디오 매칭 결과",   Inches(3.7), Inches(1.7)),
        (SHOT_SCORE,  "점수 상세 분석",        Inches(3.7), Inches(1.7)),
        (SHOT_JSONLD, "JSON-LD Schema 임베딩", Inches(3.7), Inches(1.7)),
    ]:
        img(sl, path, Inches(6.3), shot_y, width=w, height=h)
        txbox(sl, cap, Inches(6.3), shot_y + h + Inches(0.02),
              w, Inches(0.28), size=10, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER)
        shot_y += h + Inches(0.38)

    rect(sl, Inches(0.38), Inches(7.0), Inches(12.55), Inches(0.38), BG_GRAD2)
    txbox(sl, "🌐  match.yogaman.club  ·  elbee.yogaman.club  — 지금 바로 사용 가능",
          Inches(0.55), Inches(7.02), Inches(12.2), Inches(0.34),
          size=13, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — 기술 아키텍처 (matplotlib diagram embedded)
# ═══════════════════════════════════════════════════════════════════════════════

def slide_architecture():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, BG_LIGHT)
    rect(sl, 0, 0, SLIDE_W, Inches(0.12), BG_GRAD1)
    title_bar(sl, "06 · 기술 아키텍처 — 3개 파이프라인")

    # Embed the matplotlib PNG
    if os.path.exists(ARCH_PNG):
        img(sl, ARCH_PNG, Inches(0.38), Inches(1.12), width=Inches(12.55))

    # Caption strip
    rect(sl, Inches(0.38), Inches(7.05), Inches(12.55), Inches(0.35), BG_GRAD2)
    txbox(sl,
          "Pipeline A: 포즈 매칭  ·  Pipeline B: 스튜디오 매칭  ·  Pipeline C: RAG 챗봇 (Elbee)",
          Inches(0.55), Inches(7.06), Inches(12.2), Inches(0.3),
          size=11, color=BLUE_TITLE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — 기술 해자 (Moat)
# ═══════════════════════════════════════════════════════════════════════════════

def slide_moat():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, BG_LIGHT)
    rect(sl, 0, 0, SLIDE_W, Inches(0.12), BG_GRAD1)
    title_bar(sl, "07 · 기술 해자 — 복제하기 어려운 이유")

    moats = [
        ("⛔  Safety-First Kill-Switch",
         "kill_switch=TRUE 포즈를 점수 전에 제거.\n잘못 구현 시 법적 책임 → 진입 장벽.",
         RED_KILL),
        ("🏷️  E-E-A-T 강사 데이터",
         "Schema.org Person JSON-LD 자동 생성.\nGoogle SGE · Perplexity 인용 = 신뢰 자산 누적.",
         BLUE_DARK),
        ("🗺️  Haversine + 전문화",
         "거리 × 자격증 × 전문화 × 평점 동시 계산.\n단순 리스트뷰로 재현 불가.",
         EARTH),
        ("🎙️  HUD 큐잉 보이스오버",
         "Elbee 해요체 큐잉 실시간 생성.\nOllama 실패 → OpenAI 자동 전환 (99.9%).",
         LOTUS),
        ("📖  GEO 지식 그래프",
         "OCR → 키워드 인덱스 → RAG → AEO 최적화.\nAI 검색 인용은 수년간 쌓인 콘텐츠 자산.",
         SAGE),
    ]

    for i, (title, body, accent) in enumerate(moats):
        x = Inches(0.38 + (i % 3) * 4.32)
        y = Inches(1.38 + (i // 3) * 2.62)
        pill(sl, x, y, Inches(4.1), Inches(2.3), CARD_WHITE, accent, 1.3)
        txbox(sl, title, x + Inches(0.15), y + Inches(0.15),
              Inches(3.8), Inches(0.55), size=14, bold=True, color=accent, wrap=True)
        txbox(sl, body, x + Inches(0.15), y + Inches(0.7),
              Inches(3.8), Inches(1.38), size=12, color=TEXT_MID, wrap=True)

    rect(sl, Inches(0.38), Inches(7.0), Inches(12.55), Inches(0.38), BG_GRAD2)
    txbox(sl,
          "Network effect: 강사·스튜디오 데이터 증가 → 매칭 정확도 상승 → 신규 진입자 불리",
          Inches(0.55), Inches(7.02), Inches(12.2), Inches(0.34),
          size=12, italic=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — 비즈니스 모델 & 트랙션
# ═══════════════════════════════════════════════════════════════════════════════

def slide_biz():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, BG_LIGHT)
    rect(sl, 0, 0, SLIDE_W, Inches(0.12), BG_GRAD1)
    title_bar(sl, "08 · 비즈니스 모델 & 현재 트랙션")

    # Revenue streams
    txbox(sl, "매출 라인", Inches(0.38), Inches(1.35), Inches(6.0), Inches(0.42),
          size=16, bold=True, color=BLUE_TITLE)
    revenue_rows = [
        ("B2B SaaS",     "스튜디오 구독  ₩49,000/월  (리드 + 예약 API)",     BLUE_DARK),
        ("B2C 프리미엄", "사용자 맞춤 큐잉 · HUD 보이스오버  ₩9,900/월",     EARTH),
        ("API 라이선스", "Mindbody · 1club 파트너 매칭 API",                  LOTUS),
        ("AEO 컨설팅",   "강사 Schema.org 페이지 구축 대행",                  SAGE),
    ]
    for i, (t, d, accent) in enumerate(revenue_rows):
        y = Inches(1.88 + i * 0.78)
        pill(sl, Inches(0.38), y, Inches(6.0), Inches(0.65), CARD_BLUE, accent, 0.8)
        txbox(sl, t, Inches(0.52), y + Inches(0.12), Inches(1.5), Inches(0.42),
              size=12, bold=True, color=accent)
        txbox(sl, d, Inches(2.08), y + Inches(0.12), Inches(4.15), Inches(0.42),
              size=12, color=TEXT_DARK)

    # Traction
    txbox(sl, "운영 지표", Inches(7.0), Inches(1.35), Inches(6.0), Inches(0.42),
          size=16, bold=True, color=BLUE_TITLE)
    stats = [
        ("856+",  "포즈 데이터\nE-E-A-T 스키마"),
        ("8",     "서울 스튜디오\nMVP 파트너"),
        ("2",     "라이브 도메인\nCloudflare"),
        ("99.9%", "챗봇 가용성\nOllama+OpenAI"),
    ]
    for i, (n, l) in enumerate(stats):
        sx = Inches(7.0 + (i % 2) * 2.95)
        sy = Inches(1.88 + (i // 2) * 1.55)
        pill(sl, sx, sy, Inches(2.65), Inches(1.35), CARD_CREAM, BLUE_PILL, 1.0)
        txbox(sl, n, sx, sy + Inches(0.18), Inches(2.65), Inches(0.62),
              size=28, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
        txbox(sl, l, sx, sy + Inches(0.78), Inches(2.65), Inches(0.5),
              size=11, color=TEXT_MID, align=PP_ALIGN.CENTER, wrap=True)

    # GTM
    txbox(sl, "시장 진입 전략", Inches(0.38), Inches(5.12), Inches(12.5), Inches(0.42),
          size=15, bold=True, color=BLUE_TITLE)
    gtm = [
        ("01  스튜디오 확보", "요가 스튜디오 파트너십 → 리드 데이터 수집", BLUE_DARK),
        ("02  쿼리 수확",     "AI 검색 트래픽 → GEO 콘텐츠 → 사용자 유입", EARTH),
        ("03  소비자 전환",   "홈 큐잉 HUD → B2C 구독 전환",               LOTUS),
    ]
    for i, (step, desc, accent) in enumerate(gtm):
        x = Inches(0.38 + i * 4.32)
        pill(sl, x, Inches(5.62), Inches(4.1), Inches(1.2), CARD_WHITE, accent, 1.0)
        txbox(sl, step, x + Inches(0.15), Inches(5.72), Inches(3.8), Inches(0.42),
              size=13, bold=True, color=accent)
        txbox(sl, desc, x + Inches(0.15), Inches(6.12), Inches(3.8), Inches(0.6),
              size=11, color=TEXT_MID, wrap=True)

    # Funding ask
    rect(sl, Inches(0.38), Inches(6.95), Inches(12.55), Inches(0.45), BG_GRAD2, GOLD, 1)
    txbox(sl, "💰  프리시드  ·  모집 금액: ₩3억  ·  용도: 강사 DB 구축 · HUD 보이스오버 개발 · GTM",
          Inches(0.55), Inches(6.97), Inches(12.2), Inches(0.38),
          size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Closing
# ═══════════════════════════════════════════════════════════════════════════════

def slide_closing():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, BG_LIGHT)
    rect(sl, 0, 0, SLIDE_W, Inches(2.5), BG_GRAD1)
    rect(sl, 0, Inches(5.5), SLIDE_W, Inches(2.0), BG_GRAD2)
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, BLUE_DARK)

    txbox(sl, '"안전을 먼저, 효과를 그 다음에."',
          Inches(0.5), Inches(0.55), Inches(12.3), Inches(1.2),
          size=36, bold=True, italic=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
    rect(sl, Inches(3.0), Inches(1.75), Inches(7.3), Pt(2), BLUE_RULE)
    txbox(sl,
          "요가큐는 몸 상태를 이해하고, 안전한 포즈·강사·스튜디오·홈 큐잉을 연결하는 유일한 AI 매칭 플랫폼입니다.",
          Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.85),
          size=18, color=TEXT_MID, align=PP_ALIGN.CENTER, wrap=True)

    for i, (label, url) in enumerate([
        ("매칭 API",  "match.yogaman.club"),
        ("AI 챗봇",   "elbee.yogaman.club"),
        ("GitHub",    "github.com/aiegoo/aeogeo"),
    ]):
        pill(sl, Inches(1.5 + i * 3.5), Inches(3.15), Inches(3.1), Inches(0.78),
             CARD_CREAM, BLUE_PILL, 1.0)
        txbox(sl, label, Inches(1.5 + i * 3.5), Inches(3.22), Inches(3.1), Inches(0.3),
              size=12, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
        txbox(sl, url, Inches(1.5 + i * 3.5), Inches(3.52), Inches(3.1), Inches(0.3),
              size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)

    txbox(sl, "함께 AI 요가 큐레이터를 만들어요.",
          Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.75),
          size=26, bold=True, color=BLUE_TITLE, align=PP_ALIGN.CENTER)
    txbox(sl, "hello@yogaman.club  ·  AIBE5 · Team 14",
          Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.45),
          size=15, color=TEXT_MID, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════════════

slide_cover()
slide_problem()
slide_condition_check()
slide_provider_data()
slide_user_journey()
slide_demo()
slide_architecture()
slide_moat()
slide_biz()
slide_closing()

OUT = "/mnt/c/Users/hsyyu/Downloads/요가큐_발표_v3.pptx"
prs.save(OUT)
print(f"Saved → {OUT}")
print(f"Slides: {len(prs.slides)}")
